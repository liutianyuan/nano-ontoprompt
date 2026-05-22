"""
法律 & HR 测试数据 v2 生成器
增强点：
  法律 — 实体属性明细、显式实体关系图、更多触发性动作
  HR   — 显式系统自动化动作、工作流触发条件、实体属性完善
用法：cd test_data && python generate_legal_hr_v2.py
"""
import os, csv, json
from pathlib import Path

BASE = Path(__file__).parent

# ── 工具函数（同 generate_all_domains.py）──────────────────────────────────

def write_md(path, content):
    Path(path).write_text(content, encoding="utf-8")
    print(f"  ✓ {Path(path).name}")

def write_csv(path, headers, rows):
    with open(path, "w", newline="", encoding="utf-8-sig") as f:
        w = csv.writer(f); w.writerow(headers); w.writerows(rows)
    print(f"  ✓ {Path(path).name}")

def write_docx(path, title, sections):
    from docx import Document
    doc = Document()
    doc.add_heading(title, 0)
    for (heading, paras) in sections:
        doc.add_heading(heading, 1)
        for p in paras:
            if p.startswith("| "):
                rows = [r.strip().split("|")[1:-1] for r in p.strip().split("\n") if r.strip().startswith("|")]
                if rows:
                    t = doc.add_table(rows=len(rows), cols=len(rows[0]))
                    t.style = "Table Grid"
                    for i, row in enumerate(rows):
                        for j, cell in enumerate(row):
                            t.cell(i, j).text = cell.strip()
            else:
                doc.add_paragraph(p)
    doc.save(path)
    print(f"  ✓ {Path(path).name}")

def write_xlsx(path, sheets):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill
    wb = Workbook(); wb.remove(wb.active)
    for (sheet_name, headers, rows) in sheets:
        ws = wb.create_sheet(sheet_name)
        ws.append(headers)
        for cell in ws[1]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="4472C4")
        for row in rows:
            ws.append(row)
        for col in ws.columns:
            ws.column_dimensions[col[0].column_letter].width = max(
                len(str(col[0].value or "")),
                max((len(str(c.value or "")) for c in col[1:]), default=0)
            ) + 3
    wb.save(path)
    print(f"  ✓ {Path(path).name}")

def write_pptx(path, title, slides):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "业务本体知识文档 v2"
    for (stitle, bullets) in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = stitle
        tf = slide.placeholders[1].text_frame
        tf.text = bullets[0]
        for b in bullets[1:]:
            p = tf.add_paragraph(); p.text = b; p.level = 1
    prs.save(path)
    print(f"  ✓ {Path(path).name}")

def write_pdf(path, title, sections):
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    font_name = "Helvetica"
    for fp in ["C:/Windows/Fonts/simhei.ttf", "C:/Windows/Fonts/msyh.ttc"]:
        if os.path.exists(fp):
            try:
                pdfmetrics.registerFont(TTFont("CJK", fp))
                font_name = "CJK"
                break
            except Exception:
                pass
    c = canvas.Canvas(str(path), pagesize=A4)
    W, H = A4
    c.setFont(font_name, 16); c.drawString(50, H - 60, title)
    y = H - 100
    for (section_title, lines) in sections:
        if y < 100: c.showPage(); y = H - 60
        c.setFont(font_name, 13); c.drawString(50, y, section_title); y -= 24
        c.setFont(font_name, 10)
        for line in lines:
            if y < 60: c.showPage(); y = H - 60; c.setFont(font_name, 10)
            for chunk in [line[i:i+80] for i in range(0, len(line), 80)]:
                c.drawString(60, y, chunk); y -= 16
        y -= 10
    c.save()
    print(f"  ✓ {Path(path).name}")


# ══════════════════════════════════════════════════════════════════════════════
# 法律 v2
# ══════════════════════════════════════════════════════════════════════════════
def gen_legal_v2():
    D = BASE / "法律"
    print("\n[法律 v2]")
    import random; random.seed(50)

    # ── 1. legal_framework.md — 增加实体属性定义 + 显式关系 + 完整动作目录 ─────
    write_md(D / "legal_framework.md", """# 法务合规管理体系 v2

## 一、核心实体定义与属性

### 1.1 实体：合同（Contract）
| 属性 | 类型 | 说明 |
|-----|-----|-----|
| 合同编号 | String | 唯一标识，格式 CT-YYYY-NNNN |
| 合同类型 | Enum | 销售合同/采购合同/劳动合同/服务协议/战略合作/保密协议/租赁合同 |
| 合同金额 | Float | 单位：元 |
| 签署日期 | Date | 双方签字盖章日期 |
| 到期日期 | Date | 合同有效期截止 |
| 合同状态 | Enum | 草稿/审核中/有效/到期预警/已到期/争议中/已终止 |
| 适用法律 | String | 通常为"中国大陆法律" |
| 争议解决方式 | Enum | 协商/调解/仲裁/诉讼 |
| 管辖仲裁委 | String | 北京仲裁委/上海仲裁委/深圳仲裁委 |
| 法务负责人 | FK→法务人员 | 指定承办法务 |
| 业务负责人 | FK→员工 | 对接业务部门 |

### 1.2 实体：知识产权（IntellectualProperty）
| 属性 | 类型 | 说明 |
|-----|-----|-----|
| IP编号 | String | 唯一标识，格式 IP-NNNN |
| IP类型 | Enum | 发明专利/实用新型/商标/著作权/商业秘密 |
| IP名称 | String | 官方登记名称 |
| 法律状态 | Enum | 申请中/已授权/已放弃/已转让/已到期 |
| 申请日期 | Date | |
| 有效期至 | Date | 发明专利20年/实用新型10年/商标10年 |
| 保护领域 | String | 核心技术/品牌保护/商业方法/软件著作权 |
| 权利人 | FK→公司实体 | |
| 主发明人 | FK→员工 | |
| 年费状态 | Enum | 已缴/待缴/逾期 |

### 1.3 实体：争议案件（LegalDispute）
| 属性 | 类型 | 说明 |
|-----|-----|-----|
| 案件编号 | String | 格式 DIS-YYYY-NNNN |
| 争议类型 | Enum | 合同违约/货款纠纷/知识产权/劳动争议/侵权/不正当竞争 |
| 涉案金额 | Float | 单位：元 |
| 处理阶段 | Enum | 协商中/调解/仲裁/一审/二审/已结案 |
| 立案日期 | Date | |
| 承办方 | Enum | 内部法务/外部律所 |
| 外部律所 | FK→律所 | 委托律所名称 |
| 胜诉概率评估 | Enum | 高/中/低 |
| 结果 | Enum | 胜诉/败诉/和解/撤案/进行中 |
| 关联合同 | FK→合同 | 产生纠纷的合同 |

### 1.4 实体：合规检查（ComplianceAudit）
| 属性 | 类型 | 说明 |
|-----|-----|-----|
| 检查编号 | String | |
| 检查项目 | String | |
| 涉及部门 | FK→部门 | |
| 问题描述 | String | |
| 风险等级 | Enum | 高/中/低 |
| 整改期限 | Date | |
| 整改状态 | Enum | 未开始/整改中/已整改/逾期未整改 |
| 关联法规 | FK→法规 | |

### 1.5 实体：法律法规（Regulation）
| 属性 | 类型 | 说明 |
|-----|-----|-----|
| 法规编号 | String | |
| 法规名称 | String | PIPL/GDPR/劳动法/公司法等 |
| 生效日期 | Date | |
| 变更类型 | Enum | 新法规/修订/司法解释/合规指引 |
| 影响等级 | Enum | 高/中/低 |
| 应对状态 | Enum | 已评估/评估中/待评估/已整改 |

### 1.6 实体：律师/法务人员（LegalStaff）
| 属性 | 类型 | 说明 |
|-----|-----|-----|
| 工号 | String | |
| 姓名 | String | |
| 职级 | Enum | 法务专员/法务经理/法务总监/首席法律顾问 |
| 专业领域 | String | 合同法/知识产权/劳动法/数据合规 |
| 当前案件数 | Integer | |
| 执业证编号 | String | |

## 二、实体关系图谱

### 2.1 法务部门内部关系
- 首席法律顾问（GC）→ 管辖 → 法务总监
- 法务总监 → 管辖 → 法务专员（合同方向）
- 法务总监 → 管辖 → 法务专员（知识产权方向）
- 法务总监 → 管辖 → 法务专员（劳动法/合规方向）
- 法务专员 → 协作 → 外部律所

### 2.2 合同生命周期关系
- 业务部门 → 发起 → 合同草稿
- 合同草稿 → 提交审核 → 法务专员
- 法务专员 → 初审 → 合同（3个工作日内）
- 法务总监 → 终审 → 合同（重大合同）
- 合同 → 触发 → 争议案件（如违约）
- 争议案件 → 分配给 → 法务人员
- 争议案件 → 委托 → 外部律所（金额≥50万）

### 2.3 知识产权关系
- 员工 → 创造 → 知识产权成果
- 知识产权成果 → 评估 → 专利申请
- 专利申请 → 由…协助 → 专利代理机构
- 知识产权 → 保护 → 公司核心资产
- 侵权事件 → 触发 → 知识产权争议案件

### 2.4 合规管理关系
- 法律法规 → 要求遵守 → 合规检查项目
- 合规检查 → 覆盖 → 业务部门
- 合规检查 → 发现 → 合规风险
- 合规风险 → 分配整改责任 → 责任部门
- 合规风险（高）→ 上报 → 首席法律顾问

## 三、合同管理规则（扩展版）

### 3.1 合同审查权限规则
| 合同类型 | 金额阈值 | 审查人 | 审查时限 |
|--------|--------|-------|--------|
| 采购合同 | 任意金额 | 法务专员 | 3个工作日 |
| 销售合同 | > 50万 | 法务总监 | 5个工作日 |
| 战略合作协议 | 任意 | 法务总监 + 外部律所 | 10个工作日 |
| 劳动合同 | 无 | HR法务专员 | 2个工作日 |
| 知识产权转让合同 | 任意 | 法务总监 + 首席法律顾问 | 7个工作日 |

### 3.2 合同生效触发规则
- IF 合同未经法务审核 THEN 视为无效合同，禁止对外签署 → 触发动作：法务拦截通知
- IF 合同金额 > 100万 THEN 须CEO签署 → 触发动作：发起CEO审批流
- IF 合同含排他性条款 THEN 须法务总监和CEO联合审批 → 触发动作：发起联合审批流
- IF 合同内容变更超过5% THEN 须重新启动法务审核 → 触发动作：系统撤回原版本，重启流程
- IF 合同到期前90天 THEN 触发动作：发送续约提醒邮件给业务负责人+法务负责人
- IF 合同到期前30天且未续约 THEN 触发动作：升级提醒至VP，标记合同为"到期预警"
- IF 合同到期前7天且未续约 THEN 触发动作：发送紧急预警至首席法律顾问

### 3.3 合同存档规则
- IF 合同签署完成 THEN 3个工作日内 → 触发动作：上传至CMS，发送存档确认通知
- IF 合同金额 > 50万 THEN → 触发动作：同时打印纸质版存入保险柜，加密电子版

## 四、知识产权保护规则

### 4.1 商业秘密保护
- IF 员工接触核心技术信息 THEN → 触发动作：生成保密协议 + 竞业限制协议，发送签署请求
- IF 竞业补偿低于月薪30% THEN 竞业限制条款自动失效 → 触发动作：法务系统标记合规风险
- IF 离职员工加入竞争对手（竞业期内）THEN → 触发动作：停止支付竞业补偿，生成法律意见书，评估诉讼

### 4.2 专利管理规则
- IF 员工研发成果满足专利条件 THEN → 触发动作：创建专利评估工单，分配给IP法务专员
- IF 公司决定不申请专利（6个月内）THEN → 触发动作：书面通知发明人，记录放弃决定
- IF 专利年费到期前30天 THEN → 触发动作：发送年费缴纳提醒，避免专利失效
- IF 专利年费逾期未缴 THEN → 触发动作：发起紧急缴费申请，通知首席法律顾问

## 五、争议解决与法律动作

### 5.1 争议预防触发规则
- IF 合同到期后对方未付款超7天 THEN → 触发动作：法务发送书面催款函（Demand Letter）
- IF 催款超15天仍未付 THEN → 触发动作：升级为正式律师函，抄送VP和CFO

### 5.2 仲裁/诉讼启动规则
- IF 协商超30天未达成协议 THEN → 触发动作：启动仲裁申请，向仲裁委提交材料
- IF 对方资产有转移风险 THEN → 触发动作：同步申请财产保全令
- IF 涉及刑事犯罪 THEN → 触发动作：立即向公安机关报案，同步启动民事追偿

### 5.3 外部律所委托规则
- IF 争议金额 ≥ 50万 THEN → 触发动作：启动外部律所遴选，与3家以上律所比价
- IF 律所选定 THEN → 触发动作：签署委托协议，移交案件材料，设置月度汇报机制

## 六、合规管理动作

### 6.1 数据合规（PIPL/GDPR）
- IF 发生个人信息泄露 THEN → 触发动作：在72小时内向监管机构报告；通知受影响用户
- IF 用户行使删除权 THEN → 触发动作：30天内完成数据删除，书面确认发送给用户
- IF 新法规（PIPL修订）生效 THEN → 触发动作：启动影响评估，制定整改计划，安排全员培训

### 6.2 反腐败合规
- IF 礼品/招待费 > 200元/次 THEN → 触发动作：法务前置审批，留存受邀方记录
- IF 员工举报合规违规 THEN → 触发动作：启动匿名调查程序，保护举报人身份

### 6.3 合规检查整改
- IF 合规风险等级为"高" THEN → 触发动作：自动通知首席法律顾问，设置30天整改倒计时
- IF 整改逾期 THEN → 触发动作：升级至CEO报告，暂停相关业务操作
""")

    # ── 2. contract_management.docx — 增加实体关系 + 动作详情 ─────────────────
    write_docx(D / "contract_management.docx", "合同全生命周期管理操作规程 v2", [
        ("1. 实体关系：合同与相关角色", [
            "起草人（业务部门员工）→ 在CMS系统创建合同草稿 → 填写必备字段：主体、金额、标的、支付、违约、争议。",
            "法务专员 ← 系统自动分配 ← 合同草稿（按合同类型路由：采购→合同专员A，劳动→HR法务专员B）。",
            "法务总监 ← 系统推送终审任务 ← 已通过专员初审的重大合同（金额>50万 或 含知识产权）。",
            "业务VP ← 系统推送确认通知 ← 终审通过的合同，VP确认后推送给签署方。",
            "CFO/CEO ← 系统推送签署请求 ← 合同金额>100万，同时抄送审计委员会。",
        ]),
        ("2. 合同审核流程与动作触发", [
            "Step 1：起草人提交至CMS → 系统动作：验证必填字段完整性，若缺失则拒绝提交并列出缺失项。",
            "Step 2：法务专员初审（1个工作日）→ 动作：标注修改意见 / 通过。IF 实质性修改 THEN 退回起草人 + 发送修改说明。",
            "Step 3：法务总监终审（2个工作日，重大合同）→ 动作：审批通过 / 退回。IF 外部律所介入 THEN 发起外部审核工单（5个工作日）。",
            "Step 4：业务VP确认 → 动作：确认对外发送谈判稿。IF VP不认可法务意见 THEN 必须发起三方讨论会，记录会议纪要。",
            "Step 5：签署完成 → 系统动作：CMS状态变更为「有效」；触发存档流程；设置到期提醒（90/60/30/7天）。",
        ]),
        ("3. 合同履行监控自动化动作", [
            "IF 合同到期前90天 THEN 系统自动发送邮件：收件人=合同负责人+法务负责人，主题=「合同到期提醒-90天」，附合同摘要。",
            "IF 合同到期前30天且续约状态=未启动 THEN 系统自动升级通知至VP，法务标记为「高优先级续约」。",
            "IF 付款里程碑到期3天内未收款 THEN 法务发送书面催款确认；超7天 → 发送正式催款函（含违约提示）。",
            "IF 合同标的物未按时交付超15天 THEN 触发违约责任调查工单，法务评估索赔可行性，生成索赔测算报告。",
        ]),
        ("4. 争议处理标准动作序列", [
            "T+0（发现争议）：业务部门通知法务，法务建立案件档案，收集所有往来证据（邮件/微信/合同/发票）。",
            "T+3天：法务出具初步法律意见书，评估争议性质（合同违约/侵权/不当得利）和胜诉概率（高/中/低）。",
            "T+7天（争议金额<50万）：内部法务主导协商，发送和解邀请函，设定15天响应期限。",
            "T+7天（争议金额≥50万）：启动外部律所遴选（3家以上），签署委托协议，移交全部证据材料。",
            "IF 协商超30天无进展 THEN 提交仲裁申请，同步向对方发送仲裁通知，申请财产保全（若资产风险）。",
            "IF 对方提出诉讼/仲裁通知 THEN 法务须在48小时内评估并制定应对策略，启动紧急应诉流程。",
        ]),
    ])

    # ── 3. legal_database.xlsx — 增加实体关系表 ──────────────────────────────
    contract_rows = []
    types = ["销售合同","采购合同","劳动合同","服务协议","战略合作","保密协议","租赁合同","知识产权转让"]
    statuses = ["有效","到期预警","已到期","谈判中","争议中","已终止"]
    legal_staff = ["张法务（合同）","李法务（IP）","王法务（劳动）","赵总监","陈GC"]
    biz_staff = ["销售部-客户总","采购部-供应商总","产品部-产品总","运营部-运营总"]
    for i in range(1, 61):
        amt = [5, 50, 200, 1000, 5000][i % 5] * 10000
        contract_rows.append([
            f"CT-2026-{i:04d}", types[i % len(types)],
            f"甲方{i:03d}有限公司", f"乙方{i:03d}有限公司",
            amt, f"2025-{(i%12)+1:02d}-01", f"2027-{(i%12)+1:02d}-01",
            statuses[i % len(statuses)],
            "北京仲裁委" if i % 3 == 0 else "上海仲裁委",
            legal_staff[i % len(legal_staff)],
            biz_staff[i % len(biz_staff)],
        ])

    dispute_rows = []
    for i in range(1, 31):
        contract_ref = f"CT-2026-{(i*3)%60+1:04d}"
        dispute_rows.append([
            f"DIS-2026-{i:04d}",
            ["合同违约","货款纠纷","知识产权","劳动争议","侵权","不正当竞争"][i%6],
            [5, 20, 80, 350][i%4] * 10000,
            ["协商中","仲裁","诉讼","已结案"][i%4],
            f"2026-{(i%6)+1:02d}-{(i%28)+1:02d}",
            "内部法务" if i % 3 != 0 else "外部律所",
            ["金杜律师事务所","君合律师事务所","中伦律师事务所",""][i%4],
            ["高","中","低"][i%3],
            ["胜诉","败诉","和解","进行中"][i%4],
            contract_ref,
        ])

    # 实体关系表
    relation_rows = [
        ["法务部","管辖","法务总监","部门层级关系"],
        ["法务总监","管辖","法务专员-合同","岗位汇报关系"],
        ["法务总监","管辖","法务专员-IP","岗位汇报关系"],
        ["法务总监","管辖","法务专员-劳动合规","岗位汇报关系"],
        ["法务专员","审核","采购合同","合同审核关系"],
        ["法务总监","终审","重大合同(>50万)","合同审核关系"],
        ["合同","触发","争议案件","因果关系"],
        ["争议案件","分配给","法务人员","工单分配关系"],
        ["争议案件(≥50万)","委托","外部律所","外包关系"],
        ["员工","创造","知识产权成果","知识产权归属"],
        ["知识产权","评估","专利申请","流程关系"],
        ["专利申请","协助于","专利代理机构","服务关系"],
        ["合规检查","覆盖","业务部门","管理关系"],
        ["法律法规","驱动","合规检查项目","合规驱动关系"],
        ["合规风险(高)","上报","首席法律顾问","升级关系"],
        ["数据泄露事件","触发","监管机构报告","法律义务"],
        ["竞业违规","触发","仲裁/诉讼","法律动作"],
        ["合同到期","触发","续约提醒流程","自动化流程"],
    ]

    ip_rows = []
    for i in range(1, 61):
        ip_rows.append([
            f"IP-{i:04d}",
            ["发明专利","实用新型","商标","著作权","商业秘密"][i%5],
            f"知识产权-{i:03d}",
            ["已申请","已授权","申请中","已放弃","已转让"][i%5],
            f"2020-{(i%12)+1:02d}-01", f"2030-{(i%12)+1:02d}-01",
            ["核心技术","品牌保护","商业方法","软件著作权"][i%4],
            ["已缴","待缴","逾期"][i%3],
            legal_staff[i % len(legal_staff)],
        ])

    regulation_rows = []
    for i in range(1, 41):
        regulation_rows.append([
            f"REG-{i:04d}",
            ["PIPL","GDPR","劳动法","公司法","税法","反垄断法","网络安全法"][i%7],
            f"2026-{(i%6)+1:02d}-{(i%28)+1:02d}",
            ["新法规","修订","司法解释","合规指引"][i%4],
            ["高","中","低"][i%3],
            ["已评估","评估中","待评估","已整改"][i%4],
        ])

    write_xlsx(D / "legal_database.xlsx", [
        ("合同台账",
         ["合同编号","合同类型","甲方","乙方","金额(元)","签署日期","到期日期","状态","仲裁机构","法务负责人","业务负责人"],
         contract_rows),
        ("争议案件",
         ["案件编号","争议类型","涉案金额(元)","处理阶段","立案日期","承办方","外部律所","胜诉概率","结果","关联合同"],
         dispute_rows),
        ("实体关系图谱",
         ["源实体","关系类型","目标实体","关系说明"],
         relation_rows),
        ("知识产权台账",
         ["IP编号","IP类型","名称","法律状态","申请日期","有效期至","保护领域","年费状态","IP负责人"],
         ip_rows),
        ("法律法规变更",
         ["记录ID","法规名称","生效日期","变更类型","影响等级","应对状态"],
         regulation_rows),
    ])

    write_csv(D / "ip_portfolio.csv",
              ["IP编号","类型","名称","状态","申请日期","有效期至","保护领域","年费状态"],
              [[f"IP-{i:04d}",["发明专利","实用新型","商标","著作权","商业秘密"][i%5],
                f"知识产权名称{i:03d}",["已申请","已授权","申请中","已放弃"][i%4],
                f"2020-{(i%12)+1:02d}-01",f"2030-{(i%12)+1:02d}-01",
                ["核心技术","品牌保护","商业方法","软件著作权"][i%4],
                ["已缴","待缴","逾期"][i%3]] for i in range(1,61)])

    write_csv(D / "regulatory_changes.csv",
              ["记录ID","法规名称","生效日期","变更类型","影响等级","应对状态"],
              [[f"REG-{i:04d}",
                ["PIPL","GDPR","劳动法","公司法","税法","反垄断法","网络安全法"][i%7],
                f"2026-{(i%6)+1:02d}-{(i%28)+1:02d}",
                ["新法规","修订","司法解释","合规指引"][i%4],
                ["高","中","低"][i%3],
                ["已评估","评估中","待评估","已整改"][i%4]] for i in range(1,41)])

    write_pptx(D / "legal_review.pptx", "法务季度工作汇报 v2", [
        ("一、法律实体全景", [
            "核心实体：合同（186份有效）、知识产权（60项IP）、争议案件（12件在审）、合规检查（月度覆盖）、法律法规（监控40项）",
            "法务团队：GC（1）→ 法务总监（1）→ 法务专员（合同/IP/劳动各1）",
            "外部支持：金杜律师事务所（重大诉讼）、君合律师事务所（IP）、中伦律师事务所（劳动）",
        ]),
        ("二、合同管理 — 实体关系", [
            "合同→法务专员（审核关系）：专员A负责采购128份，专员B负责销售/服务58份",
            "合同→争议案件（触发关系）：6.5%的合同产生过争议（行业均值12%）",
            "合同→CFO/CEO（签署关系）：本季度CEO级别签署15份（全部>100万）",
            "关键风险：3份合同到期预警未启动续约谈判 → 已触发紧急提醒动作",
        ]),
        ("三、知识产权 — 实体状态", [
            "已授权专利：35件（发明18+实用新型17）→ 关系：全部与研发部门员工关联",
            "商标：15件已注册，3件申请中 → 关系：IP法务专员负责年费缴纳监控",
            "年费风险：2件专利年费待缴（30天内）→ 已触发：年费缴纳提醒动作",
            "侵权应对：1起竞争对手侵权 → 已触发：律师函发送、IP诉讼工单创建",
        ]),
        ("四、争议案件 — 关系与动作追踪", [
            "12件在审：8件合同违约（→关联具体合同）、4件劳动争议（→关联员工档案）",
            "外部律所委托：5件（金额>50万）→ 关系：案件→律所委托协议→月度汇报",
            "自动化动作执行情况：本季度触发「催款函」8次、「仲裁申请」2次、「财产保全」1次",
            "动作成功率：律师函发送→70%完成还款；仲裁申请→均在受理阶段",
        ]),
        ("五、合规与自动化动作统计", [
            "本季度合规检查：28次，发现高风险2项（已整改）、中风险7项（整改中）",
            "自动化动作执行：到期提醒48次、合同冻结拦截3次、数据删除请求响应5次",
            "PIPL合规：响应用户删除权申请5次（平均12天完成，目标30天内）",
            "下季度重点：完成2件竞业限制诉讼、启动5件核心技术专利申请、全员合规培训",
        ]),
    ])

    write_pdf(D / "compliance_manual.pdf", "Legal Compliance & Action Manual", [
        ("1. Contract Actions", [
            "Action: Send Renewal Reminder — Trigger: contract expires in 90 days.",
            "Action: Escalate to VP — Trigger: contract expires in 30 days, renewal not started.",
            "Action: Emergency Alert to GC — Trigger: contract expires in 7 days, no renewal.",
            "Action: Block Execution — Trigger: contract not reviewed by Legal before signing.",
            "Action: Archive to CMS — Trigger: contract signed; due within 3 business days.",
            "Action: Start CEO Approval Flow — Trigger: contract amount > 1M RMB.",
        ]),
        ("2. Dispute Actions", [
            "Action: Send Demand Letter — Trigger: payment overdue 7 days post contract term.",
            "Action: Send Lawyer Letter (律师函) — Trigger: overdue 15 days, no response.",
            "Action: File Arbitration — Trigger: negotiation exceeds 30 days with no agreement.",
            "Action: Apply Asset Preservation — Trigger: counterparty shows asset transfer risk.",
            "Action: Engage External Counsel — Trigger: dispute amount >= 500,000 RMB.",
            "Action: File Criminal Report — Trigger: criminal fraud or IP theft confirmed.",
        ]),
        ("3. IP Actions", [
            "Action: Create Patent Evaluation Ticket — Trigger: employee files R&D achievement.",
            "Action: Send Annual Fee Reminder — Trigger: patent annual fee due in 30 days.",
            "Action: Emergency Fee Filing — Trigger: annual fee overdue.",
            "Action: Send Cease-and-Desist Letter — Trigger: competitor IP infringement detected.",
            "Action: Stop Non-Compete Compensation — Trigger: ex-employee joins competitor.",
        ]),
        ("4. Compliance Actions", [
            "Action: Report to Regulator — Trigger: personal data breach detected (within 72 hrs).",
            "Action: Send Data Deletion Confirmation — Trigger: user deletion request (within 30 days).",
            "Action: Launch Compliance Training — Trigger: new regulation enacted or coverage < 95%.",
            "Action: Escalate to CEO — Trigger: high-risk compliance finding not remediated.",
            "Action: Anonymous Investigation — Trigger: whistleblower report received.",
        ]),
    ])

    print("  法律: 7个文件更新完成（v2）")


# ══════════════════════════════════════════════════════════════════════════════
# HR v2
# ══════════════════════════════════════════════════════════════════════════════
def gen_hr_v2():
    D = BASE / "HR"
    print("\n[HR v2]")
    import random; random.seed(10)

    write_md(D / "hr_policy.md", """# 人力资源管理制度 v2

## 一、核心实体定义与属性

### 1.1 实体：员工（Employee）
| 属性 | 类型 | 说明 |
|-----|-----|-----|
| 工号 | String | 唯一标识，格式 EMP-NNNN |
| 姓名 | String | |
| 职级 | Enum | P4/P5/P6/P7/P8（技术）/M1/M2/M3/M4（管理）|
| 所属部门 | FK→部门 | |
| 直属Leader | FK→员工 | |
| 当前薪资 | Float | 万元/年 |
| 市场P50薪资 | Float | 同级市场中位数 |
| 最新绩效等级 | Enum | S/A/B/C |
| 连续C季度数 | Integer | 连续绩效C的季度数 |
| 任职年限 | Float | 当前职级任职年数 |
| 晋升等待月数 | Integer | 符合晋升条件后等待月数 |
| 留存风险等级 | Enum | 高/中/低 |
| 当前状态 | Enum | 在职/试用期/PIP中/离职申请/已离职 |
| 账号状态 | Enum | 活跃/已封禁 |

### 1.2 实体：职级（Grade）
| 属性 | 类型 | 说明 |
|-----|-----|-----|
| 职级编码 | String | P4/P5/P6/P7/P8/M1-M4 |
| 薪资区间下限 | Float | |
| 薪资区间上限 | Float | |
| 最低任职年限 | Integer | 年 |
| 晋升委员会 | String | 谁负责审批该级别晋升 |
| 股票起始级别 | Boolean | P6及以上享有股票期权 |

### 1.3 实体：绩效评估（PerformanceReview）
| 属性 | 类型 | 说明 |
|-----|-----|-----|
| 评估ID | String | |
| 员工ID | FK→员工 | |
| 评估季度 | String | 2026Q1 格式 |
| OKR完成得分 | Float | 1-5分 |
| 360评估得分 | Float | 1-5分 |
| Leader评分 | Float | 1-5分 |
| 综合得分 | Float | 加权计算 |
| 绩效等级 | Enum | S/A/B/C |
| 是否纳入强制分布 | Boolean | 试用期员工除外 |

### 1.4 实体：PIP计划（PIplan）
| 属性 | 类型 | 说明 |
|-----|-----|-----|
| PIP编号 | String | |
| 员工ID | FK→员工 | |
| 触发原因 | String | 连续2季度C/单季度明显不达标 |
| 开始日期 | Date | |
| 结束日期 | Date | 通常90天 |
| 改进目标1 | String | 可量化指标 |
| 改进目标2 | String | 可量化指标 |
| 改进目标3 | String | 可量化指标 |
| 中期评估结果 | Enum | 达标/部分达标/未达标 |
| 最终结果 | Enum | 通过/未通过/进行中 |
| 法务审核状态 | Enum | 待审核/已审核 |
| HR BP负责人 | FK→HR人员 | |

### 1.5 实体：招聘需求（HiringRequest）
| 属性 | 类型 | 说明 |
|-----|-----|-----|
| 需求编号 | String | |
| 岗位名称 | String | |
| 目标职级 | FK→职级 | |
| 需求类型 | Enum | 增量HC/替补HC |
| 申请部门 | FK→部门 | |
| 审批状态 | Enum | 待审批/已审批/已关闭 |
| 开启日期 | Date | |
| 目标到岗日期 | Date | |
| 当前进度 | Enum | 简历收集/面试/Offer/入职 |
| 招聘周期(天) | Integer | 实际已用天数 |
| 超期状态 | Enum | 正常/超期7天/超期30天 |
| 招聘负责人 | FK→HR人员 | |

## 二、实体关系图谱

### 2.1 组织层级关系
- CEO → 管辖 → VP（各业务线）
- VP → 管辖 → 部门总监（M2）
- 部门总监 → 管辖 → 团队Leader（M1）
- 团队Leader → 管辖 → 工程师（P4-P7）
- HR BP → 服务于 → 各业务部门（HRBP制）
- HR BP → 协作 → 法务部（PIP文件审核）

### 2.2 员工职业发展关系
- 员工 → 属于 → 职级
- 员工 → 汇报给 → 直属Leader
- 员工 → 参与 → 绩效评估
- 绩效评估 → 触发 → 调薪决策
- 绩效评估（连续C）→ 触发 → PIP计划
- PIP计划（未通过）→ 触发 → 劝退程序
- 员工 → 申请 → 晋升评审
- 晋升评审 → 由…审批 → 晋升委员会

### 2.3 招聘流程关系
- 部门总监/VP → 发起 → 招聘需求（HC申请）
- 招聘需求 → 审批通过 → 发起招聘
- 候选人 → 通过面试 → 收到Offer
- Offer → 触发 → 背景调查
- 背景调查（通过）→ 触发 → 入职流程
- 新员工 → 完成 → 试用期评估

## 三、HR 系统自动化动作目录

### 3.1 入职自动化动作序列（Trigger: 候选人接受Offer）
1. **动作：启动背景调查** — 系统自动发送背调申请至第三方背调公司（奥创/慧安），候选人接到背调邀请
2. **动作：发送背调结果通知** — IF 背调通过 THEN 确认入职日期；IF 背调不通过 THEN 自动撤回Offer，通知HR BP
3. **动作：生成员工工号** — 入职前3个工作日，HR系统分配工号（格式：EMP-NNNN）
4. **动作：创建企业邮箱** — 入职日0:00，IT系统自动创建邮箱（姓名拼音@company.com）
5. **动作：分配系统权限** — 按职级和部门自动分配：代码仓库/OA/HR系统/项目管理工具
6. **动作：发送欢迎邮件** — 入职当天9:00，HR系统自动发送欢迎邮件（附入职手册PDF+第一周日程安排）
7. **动作：安排入职培训** — 自动注册16学时入职培训（前2天强制参与）
8. **动作：创建试用期评估表** — 入职当天，自动为直属Leader创建90天试用期评估工单

### 3.2 绩效结果自动化动作（Trigger: 季度绩效评估完成）
- IF 绩效 = S THEN：
  1. **动作：生成调薪单** — 自动计算+15%调薪，生成调薪审批单，推送给Leader→VP审批
  2. **动作：发起股票激励申请** — 创建股票激励工单，推送给薪酬委员会
  3. **动作：发送绩效反馈通知** — 向员工推送S级评定通知和奖励明细
- IF 绩效 = A THEN：
  1. **动作：生成调薪单** — 自动计算+10%调薪，生成调薪审批单
  2. **动作：发送绩效反馈通知** — 推送A级评定通知
- IF 绩效 = B THEN：
  1. **动作：生成市场薪资对比报告** — 系统拉取市场P50数据，生成薪资竞争力分析
  2. **动作：发送绩效反馈通知** — 推送B级评定通知（0-5%调薪视市场数据）
- IF 绩效 = C THEN：
  1. **动作：发送C级预警通知** — 向员工+直属Leader推送C级评定，附改进建议
  2. **动作：创建绩效改进跟踪工单** — HR BP收到工单，要求在下一季度开始前与员工完成面谈
- IF 连续2季度绩效 = C THEN：
  1. **动作：启动PIP流程** — HR系统创建PIP工单，分配给HR BP
  2. **动作：发送PIP通知书** — 生成PIP文件（含3个量化改进目标），发送给员工+法务审核
  3. **动作：设置90天倒计时** — 系统自动在30天、60天、90天发送进度检查提醒
- IF 连续3季度绩效 = C THEN：
  1. **动作：启动劝退程序** — 系统通知HR BP和法务BP，生成劝退文件草稿
  2. **动作：法律合规检查** — 自动触发劳动法合规检查（赔偿金额计算：N+1）

### 3.3 留存预警自动化动作（Trigger: 每月系统扫描）
- IF 员工薪资 < 市场P50 AND 晋升等待 > 18个月 THEN：
  1. **动作：标记高风险留存** — 员工档案状态更新为"高风险"，HR BP收到预警工单
  2. **动作：安排一对一面谈** — HR BP收到任务：7天内完成面谈，记录面谈结论
  3. **动作：生成留存方案** — HR BP填写留存方案（调薪/晋升/岗位调整），提交VP审批
- IF 关键员工（P6以上）提交离职申请 THEN：
  1. **动作：触发紧急留存评估** — HR BP在72小时内介入，评估是否反offer
  2. **动作：冻结离职申请** — 系统暂停离职流程，等待HR BP评估结论（最多72小时）
  3. **动作：生成反offer审批** — IF HR BP建议反offer THEN 自动生成调薪/晋升审批单，推送VP+CFO
- IF 核心团队（P6以上）季度离职率 > 5% THEN：
  1. **动作：上报VP和HR VP** — 系统自动生成离职率预警报告，发送给VP和HR VP
  2. **动作：启动专项留存计划** — 创建专项项目工单，由HR VP主导，30天内提交改进方案

### 3.4 招聘流程自动化动作（Trigger: 招聘需求审批通过）
- IF 招聘超过目标周期7天 THEN：
  1. **动作：发送招聘超期预警** — 通知HR BP，推送给部门VP，附当前招聘漏斗数据
- IF 招聘超过目标周期30天 THEN：
  1. **动作：升级通知VP** — 自动生成超期报告，推送给业务VP，要求7天内制定加速方案
  2. **动作：扩大招聘渠道** — 触发猎头委托评估，HR BP收到猎头渠道开启工单
- IF 候选人接受Offer THEN：
  1. **动作：启动背景调查**（见3.1）
  2. **动作：关闭招聘需求** — 系统标记该HC为"已关闭"，更新招聘漏斗数据
- IF 背调不通过 THEN：
  1. **动作：自动撤回Offer** — 系统发送撤回通知给候选人（附标准措辞）
  2. **动作：重新开启招聘需求** — 该HC状态重置为"招聘中"，重新分配至招聘负责人

### 3.5 晋升审批自动化动作
- IF 员工满足晋升条件（绩效+年限）THEN：
  1. **动作：推送晋升资格提示** — 直属Leader收到晋升资格提示，可选择是否提名
  2. **动作：创建晋升申请表** — Leader提名后，系统自动生成晋升申请表，填写晋升理由和绩效证据
- IF P4→P5晋升申请提交 THEN：**动作：创建部门总监审批工单**
- IF P5→P6晋升申请提交 THEN：**动作：创建VP+HR BP联合评审工单**
- IF P6→P7晋升申请提交 THEN：**动作：创建C-Level+外部专家评审工单**
- IF 晋升被拒绝 THEN：**动作：设置6个月冷静期标记**，员工档案更新，6个月内晋升申请自动拒绝
- IF 晋升通过 THEN：
  1. **动作：生成调薪单** — 按新职级薪资区间调整
  2. **动作：更新职级档案** — 员工档案职级字段自动更新
  3. **动作：更新系统权限** — 按新职级自动更新所有系统访问权限

### 3.6 离职自动化动作序列（Trigger: 员工提交离职申请）
1. **动作：触发留存评估**（P6以上，见3.3）
2. **动作：生成交接任务清单** — 系统自动生成标准交接任务（代码/文档/工作项），分配给员工和接替人
3. **动作：发起离职面谈预约** — HR BP收到任务，须在离职前5天完成离职面谈
4. **动作：设置账号封禁倒计时** — 最后工作日结束时（24:00），自动封禁企业邮箱/代码仓库/OA
5. **动作：生成离职证明** — 系统自动生成《离职证明》（含工作日期/职级/离职类型），HR盖章后发送
6. **动作：社保停缴申请** — HR系统在次月1日前自动发起社保停缴流程
7. **动作：检查竞业限制** — IF 员工职级 >= P6 THEN 自动触发竞业限制协议检查，法务系统生成监控任务

## 四、职级与薪酬规则

### 4.1 职级晋升规则
- IF 绩效评级 >= A 连续2个季度 AND 任职年限 >= 最低年限 THEN 可申请晋升评审
- IF 绩效评级 = S 连续1个季度 AND 直属Leader提名 THEN 可申请破格晋升（仅P4→P5）
- IF 晋升申请被拒绝 THEN 6个月冷静期，不可再次申请

### 4.2 薪酬调整规则
- IF 绩效 = S THEN 调薪15% + 股票激励
- IF 绩效 = A THEN 调薪10%
- IF 绩效 = B THEN 调薪0-5%（按市场数据）
- IF 绩效 = C THEN 无调薪
- IF 员工薪资 < 市场P25 THEN 无论绩效立即调薪（留存风险紧急处理）
- 调薪上限：单次不超过30%（特殊情况须CEO批准）
""")

    write_docx(D / "performance_management.docx", "绩效管理操作手册 v2 — 含系统动作说明", [
        ("1. 绩效评估流程与系统动作", [
            "Step 1（季度末-14天）：系统动作-发起自评通知 — 员工收到OKR填写邀请，设置7天截止时间，逾期系统自动提醒。",
            "Step 2（季度末-7天）：系统动作-发起上级评估 — Leader收到评估任务，权重60%，须在3天内完成，逾期升级通知VP。",
            "Step 3（季度末-5天）：系统动作-发起360评估 — 系统自动邀请3-5个同级/跨部门同事，匿名填写，权重20%。",
            "Step 4（季度末-2天）：系统动作-触发校准会 — 部门所有绩效汇总，HR系统检查强制分布是否合规（S≤10%，A≤40%，C≥15%）。",
            "Step 4触发：IF 分布不合规 THEN 系统动作-发送校准预警，要求Leader在24小时内调整，否则HR BP介入强制校准。",
            "Step 5（季度末）：系统动作-发送绩效结果通知 — 通知员工本季绩效等级及对应调薪/激励方案，Leader须在48小时内完成一对一反馈确认。",
        ]),
        ("2. PIP执行与系统动作", [
            "触发条件：连续2个季度绩效评级为C，或单季度OKR得分<2.5且Leader评分<3.0（明显不达标）。",
            "Step 1：系统动作-创建PIP工单 — HR BP收到工单，须在5天内与员工和Leader完成PIP启动面谈。",
            "Step 2：系统动作-生成PIP文件 — 包含3个可量化改进目标，HR BP + 法务审核（防劳动仲裁风险）。",
            "Step 3：系统动作-发送PIP通知书 — 正式通知员工（邮件+OA系统），员工须在3天内签收确认。",
            "Step 4（30天节点）：系统动作-发送中期检查提醒 — HR BP收到评估任务，员工完成中期自评。",
            "Step 5（60天节点）：系统动作-发送进度预警 — IF 中期评估=未达标 THEN 发送预警给HR VP和业务VP。",
            "Step 6（90天）：系统动作-发起最终评估 — Leader + HR BP共同评定。IF 通过 THEN 解除PIP，恢复正常流程；IF 未通过 THEN 触发劝退程序。",
            "劝退程序动作：生成劝退通知书（法务审核）、计算经济补偿金（N+1）、发起离职流程。",
        ]),
        ("3. 强制分布与校准动作", [
            "强制分布规则：S≤10%，A≤40%，B≤35%，C≥15%（部门级），各部门不可将分布压力全推到相邻部门。",
            "系统动作-自动检测：每季度评估结束后，系统扫描全公司分布，标记不合规部门。",
            "IF 部门S比例>10% THEN 系统动作-要求Leader降级处理，HR BP介入协调，须提供降级依据。",
            "IF 部门C比例<15% THEN 系统动作-标记为分布异常，HR BP约谈部门总监，须提供说明（是否有特殊业务背景）。",
            "校准会记录：全程录入HR系统，作为员工申诉证据，保存3年。",
        ]),
        ("4. 跨部门调动绩效规则", [
            "调动当季：原部门Leader评定60% + 新部门Leader评定40%，须在调动前明确告知员工。",
            "系统动作：调动申请审批通过时，HR系统自动创建当季绩效联合评估工单，分配给两位Leader。",
            "调动后第一个完整季度：新部门Leader独立评估，原部门Leader不再参与。",
            "IF 跨部门调动导致薪资低于新部门同级P50 THEN 系统动作：自动生成调薪建议，推送给新部门VP审批。",
        ]),
    ])

    emp_rows = []
    depts = ["产品研发部","销售部","客户成功部","供应链运营部","财务法务部","人力资源部"]
    grades = ["P4","P4","P5","P5","P5","P6","P6","P7","P8"]
    perf_opts = ["S","A","A","A","B","B","B","C"]
    statuses = ["在职","在职","在职","在职","PIP中","离职申请","高风险留存"]
    for i in range(1, 61):
        g = grades[i % len(grades)]
        sal = {"P4":25,"P5":40,"P6":65,"P7":105,"P8":160}[g] + random.randint(-5,5)
        market = {"P4":28,"P5":42,"P6":68,"P7":110,"P8":170}[g]
        p1, p2 = perf_opts[i%len(perf_opts)], perf_opts[(i+2)%len(perf_opts)]
        consec_c = 2 if (p1=="C" and p2=="C") else (1 if (p1=="C" or p2=="C") else 0)
        emp_rows.append([
            f"EMP{i:04d}", f"员工{i:03d}", g,
            depts[i%len(depts)], sal, market,
            p1, p2, random.randint(1,72), random.randint(0,24),
            "高" if (sal<market and random.randint(1,72)>18) else "中" if random.random()>0.7 else "低",
            consec_c, statuses[i%len(statuses)],
            "活跃" if statuses[i%len(statuses)] != "已离职" else "已封禁",
        ])

    write_xlsx(D / "employee_data.xlsx", [
        ("员工主数据",
         ["工号","姓名","职级","部门","当前薪资(万)","市场P50(万)",
          "上季绩效","本季绩效","司龄(月)","晋升等待月数",
          "留存风险","连续C季度数","当前状态","账号状态"],
         emp_rows),
        ("PIP台账",
         ["PIP编号","工号","触发原因","开始日期","目标1","目标2","目标3","中期结果","最终结果","HR BP","法务审核"],
         [
            ["PIP-2026-001","EMP0023","连续2季度C","2026-01-15","代码提交量提升50%","Bug率降至<5%","完成3个核心模块","部分达标","进行中","HR BP-张","已审核"],
            ["PIP-2026-002","EMP0041","连续2季度C","2026-02-01","销售额达Q1目标80%","客户满意度>4.0","新客户开拓>=5家","未达标","进行中","HR BP-李","已审核"],
            ["PIP-2026-003","EMP0057","单季度明显不达标","2026-03-01","文档完整率>90%","按时交付率>85%","获得3个同事正面反馈","达标","通过","HR BP-张","已审核"],
         ]),
        ("招聘漏斗",
         ["需求编号","岗位","职级","需求类型","开启日期","目标周期(天)","实际周期(天)","超期状态","简历数","面试数","Offer数","入职数","状态"],
         [
            ["HC-2026-001","AI算法工程师","P6","增量","2026-03-01",45,38,"正常",180,25,4,3,"已关闭"],
            ["HC-2026-002","产品经理","P5","替补","2026-03-15",30,42,"超期7天",220,30,5,4,"已关闭"],
            ["HC-2026-003","客户成功经理","P5","增量","2026-04-01",30,35,"超期7天",95,18,3,2,"已关闭"],
            ["HC-2026-004","供应链分析师","P5","增量","2026-04-10",30,40,"超期7天",68,12,2,1,"招聘中"],
            ["HC-2026-005","首席架构师","P7","增量","2026-04-20",60,65,"超期7天",35,8,2,0,"招聘中"],
            ["HC-2026-006","财务BP","P6","替补","2026-05-01",45,20,"正常",42,10,2,1,"招聘中"],
         ]),
        ("系统动作日志",
         ["时间","动作类型","触发条件","对象工号","执行状态","执行结果"],
         [
            ["2026-04-01 09:00","生成调薪单","绩效=S","EMP0012","已执行","调薪15%审批单发送至VP"],
            ["2026-04-01 09:00","发起股票激励申请","绩效=S","EMP0012","已执行","股票工单创建成功"],
            ["2026-04-02 08:30","启动PIP流程","连续2季度C","EMP0023","已执行","PIP工单分配给HR BP-张"],
            ["2026-04-02 08:30","发送PIP通知书","PIP启动","EMP0023","已执行","员工已签收确认"],
            ["2026-04-05 00:00","创建企业邮箱","新员工入职","EMP0061","已执行","邮箱创建成功"],
            ["2026-04-05 09:00","发送欢迎邮件","新员工入职","EMP0061","已执行","邮件发送成功"],
            ["2026-04-05 09:00","分配系统权限","新员工入职P5","EMP0061","已执行","代码仓库/OA/项目管理工具权限开通"],
            ["2026-04-10 10:00","触发留存评估","关键员工离职申请P6","EMP0045","已执行","HR BP 72小时介入任务创建"],
            ["2026-04-10 10:00","冻结离职申请","关键员工留存评估中","EMP0045","已执行","离职流程暂停，等待HR BP结论"],
            ["2026-04-15 00:00","封禁企业账号","员工离职最后工作日","EMP0038","已执行","邮箱/代码仓库/OA全部封禁"],
            ["2026-04-15 00:00","生成离职证明","员工离职","EMP0038","已执行","离职证明已发送至员工邮箱"],
            ["2026-04-15 00:00","检查竞业限制","P6员工离职","EMP0038","已执行","竞业限制监控任务已创建"],
            ["2026-04-20 09:00","发送招聘超期预警","HC-004超期7天","HC-2026-004","已执行","预警发送至HR BP和VP"],
            ["2026-05-01 09:00","发送合同到期提醒90天","无","系统","已执行","3份劳动合同到期提醒发送"],
         ]),
    ])

    perf_rows = []
    for i in range(1, 101):
        okr = round(random.uniform(2.5, 5.0), 1)
        peer = round(random.uniform(3.0, 5.0), 1)
        leader = round(random.uniform(3.0, 5.0), 1)
        final = round(okr * 0.6 + peer * 0.2 + leader * 0.2, 1)
        grade = "S" if final >= 4.5 else "A" if final >= 3.8 else "B" if final >= 3.0 else "C"
        action = {"S":"生成调薪单(+15%)+股票激励","A":"生成调薪单(+10%)","B":"生成市场薪资对比报告","C":"创建绩效改进工单+面谈任务"}[grade]
        perf_rows.append([f"EMP{i:04d}", "2026Q1", okr, peer, leader, final, grade,
                          "是" if i > 5 else "否（试用期）", action])
    write_csv(D / "performance_scores.csv",
              ["工号","季度","OKR得分","360得分","Leader评分","综合得分","绩效等级","纳入强制分布","触发系统动作"],
              perf_rows)

    risk_rows = []
    for i in range(1, 51):
        risk = ["高","高","中","中","中","低","低"][i % 7]
        action = {
            "高": "标记高风险+7天内HR BP面谈+生成留存方案提交VP",
            "中": "HR BP季度复查+薪资竞争力分析",
            "低": "年度复查"
        }[risk]
        risk_rows.append([
            f"EMP{random.randint(1,200):04d}",
            depts[i % len(depts)],
            grades[i % len(grades)],
            risk,
            random.randint(60, 130),
            random.randint(65, 135),
            random.randint(0, 36),
            ["薪酬低","晋升慢","工作压力","个人原因","竞争对手挖角"][i%5],
            action,
        ])
    write_csv(D / "retention_risk.csv",
              ["工号","部门","职级","流失风险","当前薪资","市场P50","晋升等待月数","主要风险原因","触发动作"],
              risk_rows)

    write_pptx(D / "hr_quarterly_review.pptx", "人力资源季度报告 v2 — 含动作追踪", [
        ("一、实体关系概览", [
            "员工386人 → 分布在6个部门 → 汇报给各自直属Leader → HR BP服务",
            "职级覆盖：P4(18%)、P5(42%)、P6(28%)、P7(10%)、P8(2%)",
            "实体关联：每员工关联：绩效评估记录、薪资档案、晋升历史、培训记录",
            "PIP实体：5个活跃PIP（关联员工+HR BP+法务审核状态）",
        ]),
        ("二、系统动作执行统计", [
            "本季度触发动作总计：248次",
            "调薪单生成：34次（S级9人×2动作，A级16人×1动作）",
            "PIP流程动作：5套×6步骤=30次（含通知书/中期检查/法务审核）",
            "入职自动化：8人入职×7步动作=56次（工号/邮箱/权限/欢迎邮件等）",
            "离职自动化：18人离职×5步动作=90次（含P6级别竞业检查3次）",
            "留存预警动作：23人高风险→23次面谈任务+23次留存方案提交",
        ]),
        ("三、招聘自动化动作", [
            "超期预警触发：3次（HC-002/003/004超期7天，动作：通知VP+猎头评估）",
            "背景调查：8次（均通过，0次撤回Offer）",
            "入职权限开通：8次（100%在入职当天完成）",
            "Offer撤回动作：0次（本季度无背调不通过）",
        ]),
        ("四、关键动作链路追踪", [
            "EMP0023：绩效C×2 → PIP启动 → 30天检查「部分达标」→ 60天预警发送VP → 进行中",
            "EMP0045：P6离职申请 → 冻结72h → HR BP反offer评估 → 调薪单提交VP → 留存成功",
            "EMP0038：离职通过 → 交接清单 → 离职面谈 → 账号封禁（24:00）→ 竞业检查 → 社保停缴",
            "HC-005（P7架构师）：超期65天 → 猎头渠道开启 → 仍在招聘，预计下季度到岗",
        ]),
        ("五、下季度计划", [
            "自动化升级：将「薪资竞争力分析」改为月度自动触发（当前为季度手动）",
            "PIP动作优化：增加「律师函模板」自动生成（当前需法务手工起草）",
            "晋升动作改进：P6→P7审批预计从45天缩短至30天（流程节点优化）",
            "招聘预警优化：超期3天即触发预警（当前为7天），提前干预",
        ]),
    ])

    write_pdf(D / "compensation_framework.pdf", "HR Automation & Compensation Framework", [
        ("1. Pay Philosophy & Positioning", [
            "Target: P50 for B performers, P65 for A performers, P75 for S performers.",
            "Annual salary review every April; mid-year adjustments for market correction only.",
            "Total comp = Base + Bonus + Stock Options (P6+) + Benefits.",
        ]),
        ("2. Automated Compensation Actions", [
            "Action: Generate Salary Adjustment — Trigger: performance review completed.",
            "  S grade: +15% base; Action creates approval workflow → Leader → VP → HR VP.",
            "  A grade: +10% base; Action creates approval workflow → Leader → VP.",
            "  C grade: 0% raise; Action creates coaching tracking task for HR BP.",
            "Action: Market Correction Trigger — IF salary < market P25 THEN immediate adjustment regardless of grade.",
            "Action: Bonus Pool Calculation — Triggered annually in March; system calculates pool = 15% of payroll.",
            "  IF dept KPI < 80% THEN bonus pool for dept reduced by 20% automatically.",
        ]),
        ("3. Stock Options (P6+)", [
            "Grant Schedule: P6=50k, P7=150k, P8=400k options.",
            "Vesting: 1-year cliff, then 4-year monthly vesting.",
            "Action: Grant Stock Options — Trigger: promotion to P6+ confirmed.",
            "Action: Forfeit Unvested Options — Trigger: voluntary resignation before 2-year mark.",
            "Action: Accelerate Vesting — Trigger: involuntary termination without cause (company protection).",
        ]),
        ("4. Offboarding Financial Actions", [
            "Action: Calculate Severance (N+1) — Trigger: involuntary termination after PIP failure.",
            "Action: Stop Non-Compete Compensation — Trigger: ex-employee joins competitor.",
            "Action: Final Payroll Processing — Trigger: last working day confirmed; includes unused leave payout.",
            "Action: Stop Social Insurance — Trigger: next 1st of month after last working day.",
        ]),
    ])

    print("  HR: 7个文件更新完成（v2）")


# ── 执行 ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("=" * 60)
    print("  法律 & HR 测试数据 v2 生成")
    print("=" * 60)
    gen_legal_v2()
    gen_hr_v2()
    print("\n完成！")
