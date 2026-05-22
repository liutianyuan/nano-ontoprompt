"""
多领域复杂测试数据生成器
用法：cd test_data && python generate_all_domains.py
"""
import os, csv, json
from pathlib import Path

BASE = Path(__file__).parent
DOMAINS = ["供应链", "HR", "财务", "营销", "医疗", "法律", "教育"]

for d in DOMAINS:
    (BASE / d).mkdir(exist_ok=True)

# ──────────────────────────────────────────────────────────────────────────────
# 工具函数
# ──────────────────────────────────────────────────────────────────────────────
def write_md(path, content):
    Path(path).write_text(content, encoding="utf-8")
    print(f"  ✓ {Path(path).name}")

def write_csv(path, headers, rows):
    with open(path, "w", newline="", encoding="utf-8-sig") as f:
        w = csv.writer(f); w.writerow(headers); w.writerows(rows)
    print(f"  ✓ {Path(path).name}")

def write_docx(path, title, sections):
    from docx import Document
    from docx.shared import Pt, RGBColor
    doc = Document()
    doc.add_heading(title, 0)
    for (heading, paras) in sections:
        doc.add_heading(heading, 1)
        for p in paras:
            if p.startswith("| "):  # simple table
                rows = [r.strip().split("|")[1:-1] for r in p.strip().split("\n") if r.strip().startswith("|")]
                if rows:
                    t = doc.add_table(rows=len(rows), cols=len(rows[0]))
                    t.style = "Table Grid"
                    for i, row in enumerate(rows):
                        for j, cell in enumerate(row):
                            t.cell(i, j).text = cell.strip()
            else:
                doc.add_paragraph(p)
    doc.save(path)
    print(f"  ✓ {Path(path).name}")

def write_xlsx(path, sheets):
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    wb = Workbook()
    wb.remove(wb.active)
    for (sheet_name, headers, rows) in sheets:
        ws = wb.create_sheet(sheet_name)
        ws.append(headers)
        for cell in ws[1]:
            cell.font = Font(bold=True)
            cell.fill = PatternFill("solid", fgColor="4472C4")
            cell.font = Font(bold=True, color="FFFFFF")
        for row in rows:
            ws.append(row)
        for col in ws.columns:
            ws.column_dimensions[col[0].column_letter].width = max(
                len(str(col[0].value or "")), max((len(str(c.value or "")) for c in col[1:]), default=0)
            ) + 3
    wb.save(path)
    print(f"  ✓ {Path(path).name}")

def write_pptx(path, title, slides):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    blank_layout = prs.slide_layouts[1]
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "业务本体知识文档"
    # Content slides
    for (stitle, bullets) in slides:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.title.text = stitle
        tf = slide.placeholders[1].text_frame
        tf.text = bullets[0]
        for b in bullets[1:]:
            p = tf.add_paragraph(); p.text = b; p.level = 1
    prs.save(path)
    print(f"  ✓ {Path(path).name}")

def write_pdf(path, title, sections):
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    import sys

    # Try to register a CJK font; fall back to Helvetica (ASCII only)
    font_name = "Helvetica"
    for font_path in [
        "C:/Windows/Fonts/simhei.ttf",
        "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/simsun.ttc",
    ]:
        if os.path.exists(font_path):
            try:
                pdfmetrics.registerFont(TTFont("CJK", font_path))
                font_name = "CJK"
                break
            except Exception:
                pass

    c = canvas.Canvas(str(path), pagesize=A4)
    w, h = A4
    y = h - 60
    c.setFont(font_name, 18)
    try: c.drawString(60, y, title)
    except Exception: c.drawString(60, y, title.encode("ascii","replace").decode())
    y -= 30
    c.setFont(font_name, 10)
    for (heading, body) in sections:
        if y < 80: c.showPage(); y = h - 60; c.setFont(font_name, 10)
        c.setFont(font_name, 13)
        try: c.drawString(60, y, heading)
        except Exception: c.drawString(60, y, heading.encode("ascii","replace").decode())
        y -= 18; c.setFont(font_name, 10)
        for line in body:
            for chunk in [line[i:i+80] for i in range(0, len(line), 80)]:
                if y < 80: c.showPage(); y = h - 60; c.setFont(font_name, 10)
                try: c.drawString(60, y, chunk)
                except Exception: c.drawString(60, y, chunk.encode("ascii","replace").decode())
                y -= 14
        y -= 6
    c.save()
    print(f"  ✓ {Path(path).name}")


# ══════════════════════════════════════════════════════════════════════════════
# 1. 供应链
# ══════════════════════════════════════════════════════════════════════════════
def gen_supply_chain():
    D = BASE / "供应链"
    print("\n[供应链]")

    write_md(D / "supply_chain_strategy.md", """# 供应链管理战略文件

## 1. 供应商体系

### 1.1 供应商分级
- **S级战略供应商**：年采购额>5000万，交货准时率>98%，质量合格率>99.5%。代表企业：天钢原材料有限公司、芯联电子科技、聚合包装集团。
- **A级核心供应商**：年采购额500-5000万，准时率>95%，合格率>98%。代表企业：华晟铝业、明远物流、广联五金。
- **B级普通供应商**：年采购额<500万，准时率>90%，合格率>96%。
- **C级备选供应商**：处于考察期，任意指标不达标自动降为C级，连续2个季度C级触发淘汰评审。

### 1.2 供应商评级规则
- IF 交货准时率 < 85% THEN 降一级并发送整改通知
- IF 质量合格率 < 96% AND 连续2个月 THEN 触发供应商审核
- IF 供应商评级 = C级 AND 持续2季度 THEN 启动替换流程
- IF 单次质量事故损失 > 50万元 THEN 直接暂停合作并提交管委会

## 2. 库存管理规则

### 2.1 安全库存标准
| 物料类型 | 安全库存 | 补货点 | 最大库存 |
|--------|--------|-------|--------|
| 钢材（吨） | 200 | 300 | 800 |
| 铝合金（吨） | 100 | 150 | 400 |
| 电子元件（万件） | 50 | 80 | 200 |
| 包装材料（万件） | 30 | 50 | 150 |

### 2.2 补货触发规则
- IF 库存量 < 安全库存 THEN 自动生成采购申请并通知采购员
- IF 库存量 < 安全库存*0.5 THEN 触发紧急采购，跳过常规审批直接报VP
- IF 库存量 > 最大库存*0.9 THEN 暂停补货并发送预警

## 3. 采购审批规则

| 金额区间 | 审批人 | 时限 |
|--------|-------|-----|
| < 5万元 | 采购员自审 | 1个工作日 |
| 5-50万元 | 采购部经理 | 2个工作日 |
| 50-200万元 | 供应链VP | 3个工作日 |
| > 200万元 | CFO + 董事会 | 5个工作日 |

- IF 审批超时 THEN 自动升级至上一级审批人
- IF 紧急采购标记 = True THEN 压缩审批时限至正常50%

## 4. 质量控制规则
- IF 来料检验合格率 < 98% THEN 触发全批次隔离并通知质检部
- IF 同一供应商连续3批次不合格 THEN 启动供应商专项审计
- IF 生产线不良率 > 3% THEN 停线检查，通知质量总监
- IF 客户投诉质量问题 AND 责任在供应商 THEN 启动索赔流程

## 5. 物流管理
- 标准配送时效：华东区域48小时，华南72小时，华北60小时，西部96小时
- IF 配送超时 > 24小时 THEN 自动赔付运费并降低承运商评分
- IF 货物损坏率 > 0.5% THEN 触发理赔流程并发送预警给物流总监
""")

    write_docx(D / "procurement_policy.docx", "采购管理制度", [
        ("第一章 总则", [
            "本制度适用于公司所有采购活动，包括原材料采购、辅材采购、设备采购及服务采购。",
            "采购遵循公平竞争、价值最优、合规透明的原则。所有采购活动须在ERP系统中留存完整记录。",
        ]),
        ("第二章 供应商管理", [
            "新供应商准入须通过：①资质审查 ②样品测试 ③小批量试单 ④正式合同签署四个阶段。",
            "供应商年度评估于每年12月进行，评估维度：价格竞争力（30%）、交期履约（30%）、质量水平（30%）、配合度（10%）。",
            "评估总分<70分自动降级，<60分启动替换程序，替换期间须保证备选供应商库存覆盖3个月需求。",
        ]),
        ("第三章 采购流程", [
            "需求申请：业务部门在ERP提交采购申请，填写物料编码、需求数量、需求时间、用途说明。",
            "询价比价：单次采购金额>5万须进行3家以上询价，>50万须进行正式招标。",
            "合同签署：合同须经法务审核，金额>100万须总裁签署，金额<100万采购VP签署。",
            "收货验收：仓库收货时须执行IQC入库检验，出具检验报告后方可入库。",
        ]),
        ("第四章 紧急采购", [
            "紧急采购定义：库存低于安全库存50%，或生产线面临停产风险时启动。",
            "紧急采购流程：采购员直接联系备选供应商 → 口头确认价格 → 发送采购订单 → 事后补齐审批手续。",
            "紧急采购溢价不得超过市场价15%，否则须采购VP特批。",
        ]),
    ])

    write_xlsx(D / "supplier_database.xlsx", [
        ("供应商主数据",
         ["供应商ID","供应商名称","等级","主供物料","年采购额(万)","准时率%","合格率%","联系人","状态"],
         [
            ["SUP001","天钢原材料有限公司","S","钢材",8500,98.5,99.7,"张明","有效"],
            ["SUP002","芯联电子科技","S","电子元件",6200,97.8,99.2,"李芳","有效"],
            ["SUP003","聚合包装集团","S","包装材料",3100,99.1,99.8,"王强","有效"],
            ["SUP004","华晟铝业","A","铝合金",2800,95.3,98.1,"赵雷","有效"],
            ["SUP005","明远物流","A","物流服务",1500,93.2,None,"陈梅","有效"],
            ["SUP006","广联五金","A","五金配件",900,91.5,97.8,"刘洋","有效"],
            ["SUP007","新华材料","B","辅材",450,88.0,96.2,"周杰","有效"],
            ["SUP008","东盛化工","B","化工原料",320,85.5,95.8,"吴涛","观察"],
            ["SUP009","精诚机械","B","机械部件",280,87.2,96.5,"郑华","有效"],
            ["SUP010","联达纸业","C","包装纸",150,82.1,94.3,"孙磊","整改"],
         ]),
        ("库存台账",
         ["物料编码","物料名称","规格","当前库存","安全库存","补货点","最大库存","单位","上次盘点日"],
         [
            ["MAT001","热轧钢板","Q235，5mm",185,200,300,800,"吨","2026-05-01"],
            ["MAT002","冷轧钢板","SPCC，1.5mm",312,200,300,800,"吨","2026-05-01"],
            ["MAT003","铝合金型材","6061-T6",88,100,150,400,"吨","2026-05-01"],
            ["MAT004","电子芯片","MCU-32位",42000,50000,80000,200000,"件","2026-05-01"],
            ["MAT005","电阻元件","0402，10kΩ",180000,50000,80000,200000,"件","2026-05-01"],
            ["MAT006","瓦楞纸箱","420×300×250mm",28000,30000,50000,150000,"件","2026-05-01"],
            ["MAT007","PE薄膜","0.05mm",45000,30000,50000,150000,"件","2026-05-01"],
            ["MAT008","螺栓M8","不锈钢",95000,10000,20000,50000,"件","2026-05-01"],
         ]),
        ("采购订单",
         ["订单号","供应商","物料","数量","单价(元)","总金额(万)","下单日期","要求到货","审批状态","审批人"],
         [
            ["PO-2026-0501","天钢原材料","热轧钢板",500,"4200",210,"2026-05-02","2026-05-10","已审批","供应链VP"],
            ["PO-2026-0502","芯联电子科技","MCU芯片",200000,"0.8",160,"2026-05-03","2026-05-15","已审批","供应链VP"],
            ["PO-2026-0503","聚合包装集团","瓦楞纸箱",100000,"0.32",32,"2026-05-04","2026-05-12","已审批","采购经理"],
            ["PO-2026-0504","华晟铝业","铝合金型材",200,"3800",76,"2026-05-05","2026-05-18","待审批","供应链VP"],
            ["PO-2026-0505","东盛化工","化工原料",50,"12000",60,"2026-05-06","2026-05-20","审批中","采购经理"],
         ]),
    ])

    rows_inv = []
    import random; random.seed(42)
    materials = [("MAT001","热轧钢板"),(("MAT003","铝合金型材")),(("MAT004","电子芯片")),(("MAT006","瓦楞纸箱"))]
    for i in range(1, 81):
        mat = random.choice(["MAT001","MAT003","MAT004","MAT006"])
        qty = random.randint(50, 500)
        rows_inv.append([f"2026-{random.randint(1,5):02d}-{random.randint(1,28):02d}", mat,
                         random.choice(["入库","出库","调拨","盘点"]), qty,
                         random.choice(["正常","预警","超储"]),
                         random.choice(["仓库A","仓库B","仓库C"])])
    write_csv(D / "inventory_transactions.csv",
              ["日期","物料编码","操作类型","数量","库存状态","所在仓库"], rows_inv)

    rows_logi = []
    carriers = ["顺丰","中通","申通","德邦","京东物流"]
    for i in range(1, 101):
        on_time = random.random() > 0.12
        rows_logi.append([
            f"WB-2026-{i:04d}",
            random.choice(carriers),
            f"SUP{random.randint(1,10):03d}",
            random.choice(["华东","华南","华北","西部"]),
            f"{random.randint(1,5)}天",
            "准时" if on_time else "延误",
            f"{random.uniform(0,2):.2f}%",
            f"{random.uniform(800,8000):.0f}",
        ])
    write_csv(D / "logistics_performance.csv",
              ["运单号","承运商","供应商","目的区域","实际时效","是否准时","货损率","运费(元)"], rows_logi)

    write_pptx(D / "supply_chain_review.pptx", "供应链季度运营评审", [
        ("一、核心KPI达成情况", [
            "采购准时率：96.8%（目标95%）✓",
            "库存周转天数：42天（目标45天）✓",
            "供应商质量合格率：98.7%（目标98%）✓",
            "采购成本节约率：3.2%（目标3%）✓",
        ]),
        ("二、供应商评级变动", [
            "S级：3家（天钢、芯联、聚合）—维持",
            "A级：3家（华晟、明远、广联）—维持",
            "B级升A：新华材料（连续3季度达标）",
            "C级降级：联达纸业（准时率连续2季度<85%）",
        ]),
        ("三、库存预警情况", [
            "铝合金库存88吨，低于安全库存100吨 ⚠ 已触发补货申请PO-2026-0504",
            "电子芯片库存42000件，低于安全库存50000件 ⚠ 紧急采购启动",
            "热轧钢板库存正常（185/200吨）",
        ]),
        ("四、质量事故处理", [
            "东盛化工3月批次合格率92%（低于96%阈值），已触发整改",
            "联达纸业纸箱压力测试不合格，启动全批次隔离（影响库存28000件）",
            "索赔流程已启动，预计赔付金额8万元",
        ]),
        ("五、下季度重点工作", [
            "完成2家A级供应商替换铝合金二供认证",
            "推进ERP与供应商系统直连，实现库存可视化",
            "建立供应商专项改善辅导机制（针对B级供应商）",
        ]),
    ])

    write_pdf(D / "warehouse_management.pdf", "Warehouse Management Standard", [
        ("1. Warehouse Classification", [
            "Warehouse A (Raw Materials): Steel, Aluminum — temperature 15-25C, humidity <60%",
            "Warehouse B (Electronic Components): chips, resistors — temperature 20-25C, humidity <40%",
            "Warehouse C (Finished Goods/Packaging): boxes, films — standard environment",
        ]),
        ("2. Inventory Control Rules", [
            "FIFO (First In First Out) mandatory for all materials with expiry dates.",
            "Weekly cycle count: 20% of SKUs per week, full count quarterly.",
            "IF inventory variance > 2%, mandatory root-cause investigation within 48h.",
            "Hazardous materials storage requires dedicated zone with fire suppression.",
        ]),
        ("3. Receiving Procedures", [
            "Step 1: Check PO matching (quantity, spec, supplier).",
            "Step 2: IQC sampling inspection — AQL 1.5 for critical materials.",
            "Step 3: Barcode scan and ERP update within 2 hours of receipt.",
            "IF IQC fail rate > 5%, reject entire lot and notify procurement.",
        ]),
        ("4. Dispatch Procedures", [
            "Production requests must be submitted 24h in advance.",
            "IF stock < request quantity, trigger emergency procurement alert.",
            "All dispatches require dual authorization: warehouse keeper + production supervisor.",
        ]),
    ])
    print("  供应链: 7个文件生成完成")


# ══════════════════════════════════════════════════════════════════════════════
# 2. HR
# ══════════════════════════════════════════════════════════════════════════════
def gen_hr():
    D = BASE / "HR"
    print("\n[HR]")
    import random; random.seed(10)

    write_md(D / "hr_policy.md", """# 人力资源管理制度

## 1. 职级体系

### 1.1 技术序列
| 职级 | 名称 | 薪资区间(万/年) | 最低任职年限 |
|-----|-----|--------------|-----------|
| P4 | 初级工程师 | 20-30 | 0年 |
| P5 | 中级工程师 | 30-50 | 2年 |
| P6 | 高级工程师 | 50-80 | 4年 |
| P7 | 资深工程师 | 80-130 | 6年 |
| P8 | 首席工程师 | 130-200 | 10年 |

### 1.2 管理序列（M序列）
M1（团队Leader）→ M2（部门总监）→ M3（VP）→ M4（C-Level）

## 2. 晋升规则

### 2.1 晋升触发条件
- IF 绩效评级 >= A 连续2个季度 AND 任职年限 >= 最低年限 THEN 可申请晋升评审
- IF 绩效评级 = S 连续1个季度 AND 直属Leader提名 THEN 可申请破格晋升（仅P4→P5）
- IF 晋升申请被拒绝 THEN 6个月冷静期，期间不可再次申请

### 2.2 晋升委员会
- P4→P5：部门总监审批即可
- P5→P6：VP + HR BP 联合评审
- P6→P7：C-Level + 外部专家评审
- P7→P8：CEO审批，年度名额限制<=5人/全公司

## 3. 绩效管理

### 3.1 绩效评级标准
| 等级 | 描述 | 占比约束 |
|-----|-----|--------|
| S | 超预期 | ≤10% |
| A | 达预期 | ≤40% |
| B | 基本达标 | ≤35% |
| C | 需改进 | ≥15% |

### 3.2 绩效结果应用规则
- IF 绩效 = S THEN 调薪15% + 股票激励
- IF 绩效 = A THEN 调薪10%
- IF 绩效 = B THEN 调薪0-5%（根据市场数据）
- IF 绩效 = C THEN 无调薪，启动PIP计划（绩效改进计划）
- IF 连续2季度绩效 = C THEN 触发离职风险预警，HR BP主动约谈
- IF 连续3季度绩效 = C THEN 启动劝退程序

## 4. 薪酬管理

### 4.1 薪酬调整规则
- 调薪周期：每年4月统一调薪
- IF 员工薪资 < 市场P50 AND 绩效>=B THEN 优先调薪至P50
- IF 员工薪资 < 市场P25 THEN 无论绩效立即调薪（留存风险）
- 调薪上限：单次不超过30%（特殊情况须CEO批准）

### 4.2 留存预警规则
- IF 薪资低于市场P50 AND 晋升等待>18个月 THEN 高风险留存，启动一对一面谈
- IF 关键岗位员工提交离职 THEN 72小时内HR BP介入，评估是否反offer
- IF 核心团队（P6以上）离职率 > 5%/季度 THEN 上报VP并启动专项留存计划

## 5. 招聘规则
- HC（招聘名额）审批：增量HC须VP审批；替补HC须部门总监审批
- 招聘周期目标：P4-P5不超过30天，P6不超过45天，P7以上不超过60天
- IF 招聘超期 THEN 每超7天上报HR BP，超30天上报VP
- 所有候选人须完成背景调查方可发Offer，背调不通过自动撤回Offer
""")

    write_docx(D / "performance_management.docx", "绩效管理操作手册", [
        ("1. OKR制定规范", [
            "每季度第一周完成本季度OKR制定，O不超过5个，每个O对应KR不超过4个。",
            "OKR须经直属Leader确认，P6以上须VP确认。OKR一旦锁定，季度内不得修改（特殊情况须向HR BP申请）。",
        ]),
        ("2. 绩效评估流程", [
            "Step 1: 自评（员工填写OKR完成情况，1-5分），时间：季度末最后一周",
            "Step 2: 上级评估（Leader评分，权重60%）",
            "Step 3: 360评估（3-5个同级/跨部门同事评价，权重20%）",
            "Step 4: 校准会（部门内横向拉齐，确保强制分布符合要求）",
            "Step 5: 结果反馈（Leader与员工一对一沟通，48小时内完成）",
        ]),
        ("3. PIP绩效改进计划", [
            "触发条件：连续2个季度绩效评级为C，或单季度明显不达标。",
            "PIP周期：90天，设置3个可量化的改进目标。",
            "IF PIP期间达标 THEN 恢复正常绩效流程；IF PIP结束仍未达标 THEN 启动劝退/岗位调整。",
            "PIP文件须经HR BP、法务审核后方可执行，以规避劳动仲裁风险。",
        ]),
        ("4. 特殊情况处理", [
            "产假/病假期间不参与绩效评估，按上一期绩效顺延。",
            "新员工试用期结束后首个季度仅参与自评，不纳入强制分布。",
            "跨部门调动员工当季绩效由原部门Leader评定60%+新部门Leader评定40%。",
        ]),
    ])

    emp_rows = []
    depts = ["产品研发部","销售部","客户成功部","供应链运营部","财务法务部","人力资源部"]
    grades = ["P4","P4","P5","P5","P5","P6","P6","P7","P8"]
    perf = ["S","A","A","A","B","B","B","C"]
    for i in range(1, 61):
        g = random.choice(grades)
        sal = {"P4":25,"P5":40,"P6":65,"P7":105,"P8":160}[g] + random.randint(-5,5)
        market = {"P4":28,"P5":42,"P6":68,"P7":110,"P8":170}[g]
        emp_rows.append([
            f"EMP{i:04d}", f"员工{i:03d}", g,
            random.choice(depts), sal, market,
            random.choice(perf), random.choice(perf),
            random.randint(1,72), random.choice(["在职","离职风险","PIP"])
        ])

    write_xlsx(D / "employee_data.xlsx", [
        ("员工主数据",
         ["工号","姓名","职级","所属部门","当前薪资(万)","市场P50(万)","上季绩效","本季绩效","司龄(月)","状态"],
         emp_rows),
        ("招聘漏斗",
         ["岗位","职级","开启日期","简历数","面试数","Offer数","入职数","周期(天)","状态"],
         [
            ["AI算法工程师","P6","2026-03-01",180,25,4,3,38,"已关闭"],
            ["产品经理","P5","2026-03-15",220,30,5,4,42,"已关闭"],
            ["客户成功经理","P5","2026-04-01",95,18,3,2,35,"已关闭"],
            ["供应链分析师","P5","2026-04-10",68,12,2,1,40,"招聘中"],
            ["首席架构师","P7","2026-04-20",35,8,2,0,65,"招聘中"],
            ["财务BP","P6","2026-05-01",42,10,2,1,28,"招聘中"],
         ]),
        ("培训记录",
         ["培训名称","类型","学时","必修/选修","完成率%","满意度"],
         [
            ["新员工入职培训","线下",16,"必修",98,4.5],
            ["Python数据分析","线上",8,"选修",65,4.2],
            ["OKR制定与执行","线下",4,"必修",92,4.3],
            ["供应链管理基础","线上",6,"选修",58,3.9],
            ["法律风险识别","线上",3,"必修",88,4.1],
            ["Leadership for P6+","线下",16,"必修（P6以上）",75,4.6],
         ]),
    ])

    perf_rows = []
    for i in range(1,101):
        okr = round(random.uniform(2.5,5.0),1)
        peer = round(random.uniform(3.0,5.0),1)
        final = round(okr*0.6+peer*0.2+random.uniform(3,5)*0.2,1)
        grade = "S" if final>=4.5 else "A" if final>=3.8 else "B" if final>=3.0 else "C"
        perf_rows.append([f"EMP{i:04d}",f"2026Q1",okr,peer,round(random.uniform(3,5),1),final,grade])
    write_csv(D / "performance_scores.csv",
              ["工号","季度","OKR得分","360得分","Leader评分","综合得分","绩效等级"], perf_rows)

    turnover_rows = []
    for i in range(1,51):
        risk = random.choice(["高","高","中","中","中","低","低"])
        rows2 = [f"EMP{random.randint(1,200):04d}",
                 random.choice(depts),
                 random.choice(grades),
                 risk,
                 random.randint(60,130),
                 random.randint(65,135),
                 random.randint(0,36),
                 random.choice(["薪酬低","晋升慢","工作压力","个人原因","竞争对手挖角"])]
        turnover_rows.append(rows2)
    write_csv(D / "retention_risk.csv",
              ["工号","部门","职级","流失风险","当前薪资","市场P50","等待晋升月数","主要风险原因"], turnover_rows)

    write_pptx(D / "hr_quarterly_review.pptx", "人力资源季度报告", [
        ("一、人员结构概览", ["在职总人数：386人","P6及以上占比：18%","平均司龄：2.8年","本季度新入职：24人，离职：18人"]),
        ("二、招聘达成情况", ["HC完成率：78%（39/50）","平均招聘周期：41天（目标：P5 30天 P6 45天）","校园招聘：完成春招15人","社招：AI/P7岗位难招，周期超标"]),
        ("三、绩效分布", ["S级：9%（目标≤10%）✓","A级：38%（目标≤40%）✓","B级：37%","C级：16%（目标≥15%）✓"]),
        ("四、留存预警", ["高风险员工：23人（薪资<市场P50且晋升等待>18个月）","本季度PIP启动：5人，完成改进：3人，劝退：2人","核心员工离职率：3.8%（目标<5%）✓"]),
        ("五、重点举措", ["调薪：4月完成全员年度调薪，平均涨幅8.5%","职级校准：P5→P6晋升18人，P6→P7晋升3人","启动P7以上薪酬竞争力专项调研"]),
    ])

    write_pdf(D / "compensation_framework.pdf", "Compensation & Benefits Framework", [
        ("1. Pay Philosophy", [
            "Target market positioning: P50 for B performers, P65 for A performers, P75 for S performers.",
            "Annual salary review in April based on merit and market data.",
            "Total compensation = Base Salary + Performance Bonus + Stock Options (P6+).",
        ]),
        ("2. Bonus Structure", [
            "Annual bonus pool: 15% of total payroll, distributed by grade and performance.",
            "S grade: 30% of annual base; A grade: 20%; B grade: 10%; C grade: 0%.",
            "IF department KPI achievement < 80% THEN bonus pool reduced by 20%.",
        ]),
        ("3. Stock Option Plan", [
            "Eligible: P6 and above, post 1-year cliff, 4-year vesting schedule.",
            "Grant size: P6 = 50k options, P7 = 150k options, P8 = 400k options.",
            "IF employee voluntarily leaves before 2-year mark, unvested options forfeited.",
        ]),
        ("4. Benefits", [
            "Medical insurance: employee + spouse + 2 children covered.",
            "Annual leave: P4-P5 = 10 days, P6-P7 = 15 days, P8+ = 20 days.",
            "Education allowance: up to 10,000 RMB/year for job-related training.",
        ]),
    ])
    print("  HR: 7个文件生成完成")


# ══════════════════════════════════════════════════════════════════════════════
# 3. 财务
# ══════════════════════════════════════════════════════════════════════════════
def gen_finance():
    D = BASE / "财务"
    print("\n[财务]")
    import random; random.seed(20)

    write_md(D / "financial_controls.md", """# 财务内部控制制度

## 1. 采购与付款控制

### 1.1 采购审批权限
| 金额区间 | 审批人 | 系统控制 |
|--------|-------|--------|
| < 1万元 | 部门经理 | ERP自动放行 |
| 1-10万元 | 财务BP + 部门VP | 需双人审批 |
| 10-100万元 | CFO | 需附合同扫描件 |
| > 100万元 | CFO + CEO + 审计委员会 | 需董事会决议 |

### 1.2 付款控制规则
- IF 供应商发票金额 ≠ 采购订单金额（差异>0.5%）THEN 付款冻结，触发差异调查
- IF 供应商银行账户变更 THEN 须经双人验证并重新录入，生效后24小时内通知CFO
- IF 同一供应商单月付款 > 500万 THEN 自动发送CFO预警邮件

## 2. 费用报销控制

### 2.1 报销限额标准
| 费用类型 | 限额 | 超限处理 |
|--------|-----|--------|
| 国内差旅住宿 | 600元/晚 | 需VP审批 |
| 国际出差住宿 | 200USD/晚 | 需CFO审批 |
| 客户招待 | 500元/人次 | 需提供客户名单 |
| 交通（打车） | 500元/次 | 需附行程说明 |
| 礼品 | 200元/份 | 禁止现金礼品 |

### 2.2 报销审批规则
- IF 报销金额 > 10000元 THEN 须附原始发票（不接受复印件）
- IF 同一员工单月报销 > 20000元 THEN 触发财务BP审查
- IF 报销事项 = 礼品/娱乐 THEN 须填写反腐合规声明

## 3. 预算管理

### 3.1 预算控制规则
- IF 部门实际支出 > 预算90% THEN 发送预警通知给部门VP
- IF 部门实际支出 > 预算100% THEN 冻结该科目，须CFO特批方可继续
- IF 单笔支出 > 剩余季度预算20% THEN 须提前申请预算调整

### 3.2 预算调整规则
- 季度内预算调整须CFO审批，累计调整不超过原预算10%
- IF 全年预算调整 > 10% THEN 须董事会审批
- 人员预算（HC）调整须HR VP + CFO联合审批

## 4. 应收账款管理

### 4.1 账龄预警规则
| 账龄 | 处理规则 |
|-----|--------|
| 0-30天 | 正常跟进 |
| 31-60天 | 财务发送催收提醒 |
| 61-90天 | 客户成功介入，升级为VP处理 |
| 91-180天 | 法务发送律师函，暂停新增授信 |
| >180天 | 计提50%坏账准备，启动诉讼评估 |
| >365天 | 计提100%坏账，申请核销 |

### 4.2 坏账计提规则
- IF 账龄91-180天 THEN 按余额50%计提
- IF 账龄>180天 THEN 按余额100%计提
- IF 客户宣告破产 THEN 立即100%计提，无论账龄
- IF 计提金额 > 100万 THEN 须CFO审批并在财报中披露

## 5. 资产折旧规则
- 固定资产折旧：IT设备3年直线法，机械设备10年，房屋建筑40年
- IF 资产残值 < 账面价值5% THEN 加速折旧直至完全摊销
- IF 资产闲置>6个月 THEN 发起资产处置评估
""")

    write_docx(D / "month_end_close.docx", "月末结账操作手册", [
        ("1. 结账时间表", [
            "D+1（次月1日）：完成所有报销单录入，截止时间18:00",
            "D+2：财务BP完成收入确认，AR团队完成发票核对",
            "D+3：完成固定资产折旧计提，完成预提费用确认",
            "D+4：完成成本结转，生成初步利润表",
            "D+5：CFO Review，处理异常科目，发布最终报表",
        ]),
        ("2. 关键控制点", [
            "收入确认：所有收入须有已签署合同+交付确认函，缺一不可。",
            "IF 收入确认缺少合同支撑 THEN 该收入挂应收，不得确认当月收入。",
            "成本配比：收入与成本须在同一期间确认，禁止跨期调整。",
            "银行对账：财务必须每月核对所有银行账户，差异须在结账前消除。",
        ]),
        ("3. 异常处理规则", [
            "IF 科目余额异常（如资产出现负值、费用出现贷方余额）THEN 当天内查明原因并处理。",
            "IF 供应商对账差异 > 1000元 THEN 不得关闭该供应商对账，须在下月5日前处理。",
            "IF 税务申报数据与账面不符 THEN 立即上报税务BP，不得擅自处理。",
        ]),
        ("4. 税务合规", [
            "增值税申报：每月15日前完成上月增值税申报及缴纳。",
            "企业所得税：季度预缴，次月15日前完成；年度汇算清缴，次年5月31日前完成。",
            "IF 税务稽查 THEN 所有原始凭证须10年内保存，不得销毁。",
        ]),
    ])

    months = ["2026-01","2026-02","2026-03","2026-04","2026-05"]
    depts2 = ["产品研发","销售","客户成功","供应链","财务法务","人力资源","市场营销","行政"]
    budget_rows = []
    for d in depts2:
        budget = random.randint(200,2000)
        for m in months:
            actual = budget//5 * random.uniform(0.7,1.3)
            budget_rows.append([d, m, budget//5, round(actual,1), round(actual/(budget//5)*100,1),
                                 "正常" if actual<budget//5*1.0 else "超支" if actual>budget//5 else "正常"])

    ar_rows = []
    for i in range(1,61):
        age = random.choice([15,25,45,75,120,200,400])
        amt = random.randint(5,500)*10000
        prov = 0 if age<=90 else amt*0.5 if age<=180 else amt
        ar_rows.append([f"CUS{i:04d}", f"客户{i:03d}",
                        f"{random.randint(2025,2026)}-{random.randint(1,12):02d}-{random.randint(1,28):02d}",
                        amt, age, round(prov,0), "正常" if age<=60 else "预警" if age<=90 else "催收"])

    write_xlsx(D / "financial_data.xlsx", [
        ("预算执行",
         ["部门","月份","预算(万)","实际(万)","完成率%","状态"], budget_rows),
        ("应收账款",
         ["客户ID","客户名称","发票日期","金额(元)","账龄(天)","已计提坏账","状态"], ar_rows),
        ("科目余额",
         ["科目编码","科目名称","期初余额","本期借方","本期贷方","期末余额","说明"],
         [
            ["1001","库存现金",50000,120000,95000,75000,"正常"],
            ["1002","银行存款",8500000,15000000,12000000,11500000,"正常"],
            ["1122","应收账款",3200000,5000000,4200000,4000000,"含逾期360万"],
            ["2202","应付账款",2100000,1800000,2500000,2800000,"正常"],
            ["4001","主营业务收入",0,0,12500000,12500000,"当月确认"],
            ["6401","主营业务成本",0,8200000,0,8200000,"毛利率34.4%"],
         ]),
    ])

    expense_rows = []
    types = ["差旅","客户招待","办公用品","技术服务","培训"]
    for i in range(1,121):
        limit = {"差旅":600,"客户招待":500,"办公用品":200,"技术服务":10000,"培训":2000}
        t = random.choice(types)
        amt = random.randint(100,int(limit[t]*1.5))
        over = amt > limit[t]
        expense_rows.append([f"EXP{i:05d}", f"EMP{random.randint(1,60):04d}",
                             t, amt, "是" if over else "否",
                             f"2026-{random.randint(1,5):02d}-{random.randint(1,28):02d}",
                             "待审核" if over else "已通过"])
    write_csv(D / "expense_reports.csv",
              ["报销单号","申请人","费用类型","金额(元)","是否超标","申请日期","审批状态"], expense_rows)

    cashflow_rows = []
    for i in range(1,61):
        cashflow_rows.append([f"2026-{random.randint(1,5):02d}-{random.randint(1,28):02d}",
                              random.choice(["经营","投资","筹资"]),
                              random.choice(["收款","付款"]),
                              random.randint(10000,5000000),
                              random.choice(["客户收款","供应商付款","工资","税款","贷款还款"]),
                              f"BNK{random.randint(1,3):01d}"])
    write_csv(D / "cash_flow.csv",
              ["日期","活动类型","收付方向","金额(元)","摘要","银行账户"], cashflow_rows)

    write_pptx(D / "cfo_review.pptx", "CFO月度财务评审", [
        ("一、收入与利润", ["当月收入：1,250万（同比+18%）✓","毛利率：34.4%（目标32%）✓","净利润：210万（利润率16.8%）","费用率：17.6%（目标<20%）✓"]),
        ("二、现金流", ["期末现金余额：1,150万（安全线：500万）✓","经营活动现金净流入：320万","应收账款周转天数：38天（目标45天）✓"]),
        ("三、预算超支预警", ["市场营销部：超预算12%，已冻结；需CFO特批","产品研发部：预算使用率88%，正常","销售部：预算使用率95%，预警"]),
        ("四、应收账款风险", ["90天以上逾期：360万（含3家高风险客户）","已启动法务催收：2家，金额120万","本月计提坏账准备：55万"]),
        ("五、下月重点", ["完成年度税务汇算：5月31日截止","审计配合：外部审计进场5月20日","推进ERP升级：财务模块上线测试"]),
    ])

    write_pdf(D / "audit_report.pdf", "Internal Audit Report Q1 2026", [
        ("Executive Summary", [
            "Audit period: January 1 - March 31, 2026.",
            "Scope: Procurement, Expense Reimbursement, Revenue Recognition.",
            "Overall assessment: Satisfactory with 3 medium-risk findings.",
        ]),
        ("Finding 1: Procurement Approval Bypass", [
            "2 instances where purchases >100k RMB approved by VP only (should require CFO).",
            "Risk: Authorization control weakness.",
            "Recommendation: Implement system hard-stop in ERP for amounts exceeding threshold.",
            "Management response: ERP control to be implemented by June 30, 2026.",
        ]),
        ("Finding 2: Expense Report Timeliness", [
            "15% of expense reports submitted >30 days after the expense date.",
            "Risk: Period cutoff misstatement.",
            "Recommendation: Enforce 15-day submission policy via system reminder.",
        ]),
        ("Finding 3: AR Follow-up", [
            "3 customers with balances >180 days not yet referred to legal (total RMB 850k).",
            "Recommendation: Automate escalation trigger at 180-day mark.",
        ]),
    ])
    print("  财务: 7个文件生成完成")


# ══════════════════════════════════════════════════════════════════════════════
# 4. 营销
# ══════════════════════════════════════════════════════════════════════════════
def gen_marketing():
    D = BASE / "营销"
    print("\n[营销]")
    import random; random.seed(30)

    write_md(D / "marketing_strategy.md", """# 营销与客户运营战略

## 1. 客户分层体系

### 1.1 客户分级标准
| 等级 | 名称 | ARR标准 | 健康度要求 | CSM配置 |
|-----|-----|--------|---------|--------|
| S级 | 战略客户 | > 500万/年 | ≥75分 | 专属CSM（1:5） |
| A级 | 重点客户 | 100-500万/年 | ≥65分 | 核心CSM（1:20） |
| B级 | 成长客户 | 20-100万/年 | ≥60分 | 标准CSM（1:50） |
| C级 | 长尾客户 | < 20万/年 | 无要求 | 数字化服务 |

### 1.2 客户分级规则
- IF ARR > 500万 THEN 自动升为S级，分配专属CSM
- IF ARR连续2季度下降>20% THEN 降一级并触发CSM介入
- IF 客户主动申请降配 AND ARR达标 THEN 保留60天观察期

## 2. 健康度管理

### 2.1 健康度评分模型
| 维度 | 权重 | 评分标准 |
|-----|-----|--------|
| 产品使用率 | 35% | MAU/授权用户数 |
| NPS净推荐值 | 25% | 最近一次NPS调研 |
| 支持工单 | 20% | 工单解决率及满意度 |
| 商务关系 | 20% | 联系人稳定性、付款及时性 |

### 2.2 健康度预警规则
- IF 健康度分 < 60 THEN 红色预警，CSM须48小时内联系客户
- IF 健康度分 60-70 THEN 黄色预警，CSM须本周内制定改善计划
- IF 健康度连续3个月下降 THEN 升级至CSM Leader处理
- IF 关键联系人（决策人）离职 THEN 触发关系重建流程，72小时内介入

## 3. 续约管理

### 3.1 续约触发规则
- IF 合同到期前90天 AND 健康度 < 70 THEN 启动高风险续约流程（CSM + AE共同跟进）
- IF 合同到期前60天 THEN 正常续约流程（CSM推进报价）
- IF 合同到期前30天未完成续约 THEN 上报VP，启动应急挽留
- IF 续约报价 = 原价*1.2 AND 客户健康度 < 65 THEN 不建议涨价，须CSM Leader审批

### 3.2 流失风险规则
- IF NPS < 20 THEN 发送高风险预警至CSM + VP
- IF 产品MAU/授权用户 < 30% 连续2个月 THEN 触发使用激活计划
- IF 客户提交取消申请 THEN 24小时内CSM主管介入，评估是否可挽回

## 4. 线索管理

### 4.1 线索评分模型（MQL触发条件：80分以上）
- 行为分（60%）：访问官网(+5/次)、下载白皮书(+15)、注册试用(+30)、参加Webinar(+20)
- 属性分（40%）：企业规模>500人(+20)、决策人联系方式(+15)、目标行业(+10)

### 4.2 线索分配规则
- IF 线索评分 >= 80 THEN 自动成为MQL，分配给对应区域销售
- IF 线索来自S级客户域名 THEN 直接分配给S级销售团队，跳过评分流程
- IF 销售72小时内未跟进MQL THEN 线索归还线索池并发送预警给销售经理

## 5. 渠道投放规则
- IF SEM单渠道CAC > LTV*0.3 THEN 削减该渠道预算50%，转至内容营销
- IF 某渠道月获客MQL < 10 THEN 下月暂停该渠道投放
- IF 新客户首单ARR < 20万 THEN 划入C级数字化服务，不分配CSM
""")

    write_docx(D / "csm_playbook.docx", "客户成功运营手册", [
        ("1. 客户入职（Onboarding）", [
            "合同签署后3个工作日内，CSM完成首次欢迎拜访（视频会议或现场）。",
            "制定客户成功计划（Success Plan），包含：目标KPI、里程碑时间表、风险点识别。",
            "IF 客户30天内未完成初始配置 THEN 触发技术支持介入，并在健康度中扣分。",
            "60天Check-in：评估产品使用率，如MAU<预期50%，启动使用激活专项。",
        ]),
        ("2. 日常运营节奏", [
            "S级客户：每月现场EBR（业务评审），每周微信/邮件简报",
            "A级客户：每季度视频EBR，每月数据报告",
            "B级客户：每半年健康度评估，自助客户门户",
            "IF 任何等级客户健康度<60 THEN 立即升级为'救援模式'，周频跟进",
        ]),
        ("3. 续约流程", [
            "T-90天：发起内部续约评审会，评估健康度、扩展机会、风险点",
            "T-60天：向客户发起续约沟通，递交续约提案",
            "T-30天：如未签约，上报VP，启动挽留方案（可能涉及折扣）",
            "续约折扣审批：≤5% CSM Leader批准，5-15% VP批准，>15% CEO批准",
        ]),
        ("4. 扩展销售（Upsell/Crosssell）", [
            "客户健康度>75且满6个月使用后，CSM主动发起扩展销售对话",
            "IF 扩展金额 > 50万 THEN CSM + AE联合跟进",
            "每季度CSM须识别至少1个扩展机会（A/S级客户必须完成）",
        ]),
    ])

    cust_rows = []
    for i in range(1,81):
        arr = random.choice([random.randint(5,18), random.randint(20,95), random.randint(100,490), random.randint(500,2000)])*10000
        grade = "S" if arr>5000000 else "A" if arr>1000000 else "B" if arr>200000 else "C"
        health = random.randint(45,95)
        nps = random.randint(-20,70)
        mau = random.uniform(0.2,0.95)
        cust_rows.append([f"CUS{i:04d}", f"客户{i:03d}企业", grade, arr, health,
                          nps, f"{mau:.0%}", random.choice(["正常","黄色预警","红色预警"]),
                          f"2026-{random.randint(3,12):02d}-01"])

    campaign_rows = []
    channels = ["SEM","内容营销","社交媒体","邮件营销","线下活动","合作伙伴"]
    for i in range(1,61):
        ch = random.choice(channels)
        spend = random.randint(5,500)*1000
        mql = random.randint(1,80)
        won = random.randint(0,max(1,mql//5))
        arr = won * random.randint(200000,1500000)
        cac = spend/max(won,1)
        roi = arr/spend if spend else 0
        campaign_rows.append([f"2026-{random.randint(1,5):02d}", ch, spend, mql, won, arr, round(cac,0), round(roi,2)])

    write_xlsx(D / "marketing_data.xlsx", [
        ("客户健康度",
         ["客户ID","客户名称","等级","ARR(元)","健康度","NPS","MAU达成率","预警状态","合同到期"],
         cust_rows),
        ("渠道投放",
         ["月份","渠道","投入(元)","MQL数","赢单数","带来ARR(元)","CAC(元)","ROI"],
         campaign_rows),
        ("线索漏斗",
         ["阶段","本月数量","上月数量","转化率%","平均停留天数"],
         [
            ["官网访客",8500,7200,None,None],
            ["注册用户",620,540,"7.3%",None],
            ["MQL",180,155,"29%",3],
            ["SQL",72,60,"40%",7],
            ["提案",35,28,"49%",14],
            ["赢单",18,14,"51%",21],
         ]),
    ])

    lead_rows = []
    sources = ["SEM","内容","社媒","活动","转介绍"]
    for i in range(1,101):
        score = random.randint(20,100)
        lead_rows.append([f"LEAD{i:05d}", random.choice(sources),
                          f"2026-{random.randint(1,5):02d}-{random.randint(1,28):02d}",
                          score, "MQL" if score>=80 else "潜在",
                          random.randint(50,5000), random.choice(["技术","金融","制造","零售","医疗"]),
                          random.choice(["已联系","未跟进","无效","培育中"])])
    write_csv(D / "lead_data.csv",
              ["线索ID","来源渠道","创建日期","评分","线索状态","企业规模","行业","跟进状态"], lead_rows)

    nps_rows = []
    for i in range(1,61):
        score = random.randint(-50,100)
        nps_rows.append([f"CUS{random.randint(1,80):04d}",
                         f"2026-{random.randint(1,5):02d}",
                         score,
                         "推荐者" if score>=50 else "被动者" if score>=0 else "批评者",
                         random.choice(["产品功能强","服务响应慢","价格偏高","易用性好","技术支持好","功能缺失"])])
    write_csv(D / "nps_survey.csv",
              ["客户ID","调研月份","NPS分值","客户类型","主要反馈"], nps_rows)

    write_pptx(D / "marketing_review.pptx", "营销季度业务评审", [
        ("一、获客情况", ["本季MQL：547个（目标500）✓","赢单率：25%（目标22%）✓","新增ARR：2,340万","平均成交周期：42天"]),
        ("二、客户健康度", ["S级客户健康度平均：81分（目标≥75）✓","A级：72分（目标≥65）✓","红色预警客户：8家（需立即介入）","NPS：42（目标≥40）✓"]),
        ("三、续约情况", ["Q1续约率：91%（目标90%）✓","流失ARR：350万（含2家S级客户）","高风险续约（T-90天内）：12家"]),
        ("四、渠道效率", ["最优渠道：转介绍（ROI 8.5x）","表现最差：线下活动（CAC 15万，超标）","SEM CAC：3.2万，接近阈值，预警"]),
        ("五、下季度计划", ["启动S级客户专项扩展计划（目标增扩ARR 800万）","下线低效线下活动，预算转至内容营销","新增行业解决方案白皮书3篇"]),
    ])

    write_pdf(D / "crm_automation_rules.pdf", "CRM Automation Rules Specification", [
        ("1. Lead Scoring Rules", [
            "Website visit: +5 points per session (max 20 points).",
            "Whitepaper download: +15 points.",
            "Free trial registration: +30 points.",
            "Webinar attendance: +20 points.",
            "IF lead score >= 80 THEN convert to MQL and assign to sales rep.",
            "IF lead score >= 80 AND company size > 500 THEN assign to enterprise sales.",
        ]),
        ("2. Customer Health Monitoring", [
            "Health score calculated weekly: Product Usage 35% + NPS 25% + Support 20% + Relationship 20%.",
            "IF health score drops > 10 points in one week THEN send alert to CSM.",
            "IF health score < 60 THEN create urgent task for CSM, due within 48 hours.",
            "IF health score < 50 THEN escalate to CSM Manager and VP.",
        ]),
        ("3. Renewal Automation", [
            "T-90: Auto-create renewal opportunity in CRM, assign to CSM.",
            "T-60: Auto-send renewal proposal template.",
            "T-30: Auto-escalate to CSM Manager if deal not closed.",
            "IF renewal not signed by T-0 THEN mark as churn, initiate win-back sequence.",
        ]),
    ])
    print("  营销: 7个文件生成完成")


# ══════════════════════════════════════════════════════════════════════════════
# 5. 医疗
# ══════════════════════════════════════════════════════════════════════════════
def gen_medical():
    D = BASE / "医疗"
    print("\n[医疗]")
    import random; random.seed(40)

    write_md(D / "clinical_protocols.md", """# 临床诊疗规范手册

## 1. 高血压诊疗规范

### 1.1 诊断标准
- 收缩压（SBP）≥140 mmHg 且/或 舒张压（DBP）≥90 mmHg（非同日三次测量）
- IF SBP ≥ 180 OR DBP ≥ 120 THEN 高血压危象，立即急诊处理
- IF SBP 130-139 AND DBP 80-89 THEN 高血压前期，生活方式干预，每3个月随访

### 1.2 治疗方案选择规则
| 分级 | SBP | 一线治疗 |
|-----|-----|--------|
| 1级 | 140-159 | 生活方式干预4周；无效加用ACEI/ARB |
| 2级 | 160-179 | 立即启动ACEI/ARB + CCB联合治疗 |
| 3级 | ≥180 | 住院治疗，静脉降压药 |

### 1.3 用药禁忌规则
- IF 妊娠 THEN 禁用ACEI、ARB（致畸）；使用甲基多巴或硝苯地平缓释片
- IF 慢性肾病eGFR < 30 mL/min THEN 禁用噻嗪类利尿剂；慎用ACEI
- IF 哮喘或COPD THEN 禁用β受体阻滞剂（美托洛尔/比索洛尔）
- IF 高血钾（K+ > 5.5 mEq/L）THEN 禁用ACEI/ARB，改用CCB

## 2. 2型糖尿病诊疗规范

### 2.1 诊断标准
- 空腹血糖（FPG）≥ 7.0 mmol/L，或OGTT 2小时血糖 ≥ 11.1 mmol/L
- HbA1c ≥ 6.5%（经标准化实验室检测）
- IF FPG 6.1-6.9 THEN 空腹血糖受损（IFG），启动糖尿病风险管理

### 2.2 血糖控制目标
| 人群 | HbA1c目标 | 空腹血糖 | 餐后2h血糖 |
|-----|---------|--------|---------|
| 一般成人 | < 7.0% | 4.4-7.0 | < 10.0 |
| 老年人（>70岁） | < 8.0% | 5.0-8.0 | < 12.0 |
| 孕期 | < 6.0% | 3.3-5.3 | < 6.7 |

### 2.3 用药调整规则
- 一线：二甲双胍（IF eGFR ≥ 45 mL/min）
- IF eGFR 30-44 THEN 减量二甲双胍；IF eGFR < 30 THEN 停用二甲双胍，改用SGLT-2抑制剂
- IF HbA1c > 9% 且已用二甲双胍最大剂量 THEN 启动胰岛素治疗
- IF 低血糖（血糖 < 3.9 mmol/L）THEN 立即口服15g快速碳水，15分钟后复测

### 2.4 并发症筛查规则
- IF 糖尿病确诊 THEN 每年筛查：肾功能（eGFR+尿微量白蛋白）、眼底、神经病变、足部
- IF 尿微量白蛋白/肌酐 > 30 mg/g THEN 诊断糖尿病肾病早期，启动ACEI/ARB保护肾功能
- IF 视网膜病变 = 增殖期 THEN 转眼科，评估激光治疗或玻璃体手术

## 3. 药物安全规则

### 3.1 药物相互作用
- 华法林 + 阿司匹林 → 出血风险增加，须INR监测（目标2.0-3.0）
- 他汀类 + 贝特类 → 横纹肌溶解风险，避免联用；必须联用时剂量减半并监测CK
- 地高辛 + 胺碘酮 → 地高辛血药浓度升高，须减量50%并监测

### 3.2 老年患者特殊规则
- IF 年龄 > 75岁 THEN 所有药物起始剂量不超过标准剂量50%，缓慢滴定
- IF 多药联用（≥5种）THEN 进行用药重整（Medication Reconciliation）
- IF 肾功能eGFR < 60 THEN 对经肾排泄药物进行剂量调整
""")

    write_docx(D / "treatment_procedures.docx", "诊疗操作规程", [
        ("1. 急诊接诊流程", [
            "预检分诊：护士在患者到达5分钟内完成生命体征评估（血压、心率、体温、SpO2）",
            "IF SpO2 < 90% THEN 立即给氧，开启急救绿色通道",
            "IF 收缩压 < 90 mmHg THEN 触发休克处理流程，开通静脉通路，通知值班医生",
            "胸痛患者：10分钟内完成12导联心电图，30分钟内完成肌钙蛋白检测",
            "IF 心电图提示ST抬高 THEN 立即启动胸痛中心流程，30分钟内完成导管室准备（Door-to-Balloon < 90分钟）",
        ]),
        ("2. 用药安全核查", [
            "医嘱双核查制度：护士执行前须核对医嘱（药品名、剂量、途径、时间），由另一名护士交叉核查",
            "高警示药品（胰岛素、氯化钾、肝素等）：须两名护士共同核查，独立计算剂量",
            "IF 静脉推注速度超过规定限速 THEN 护理系统自动报警，停止推注",
        ]),
        ("3. 感染控制规程", [
            "手卫生：所有接触患者前后须使用免洗手消毒液或肥皂洗手（七步法，20-30秒）",
            "IF 患者诊断为多重耐药菌（MRSA/VRE/CRE）THEN 启动隔离预防措施：单间隔离/接触隔离，进入须穿戴手套+隔离衣",
            "IF 针刺伤 THEN 立即执行职业暴露处理流程：伤口冲洗→报告→评估→预防用药（30分钟内）",
        ]),
        ("4. 随访管理规程", [
            "慢病患者（高血压/糖尿病/冠心病）：至少每3个月随访一次",
            "IF HbA1c > 8% THEN 增加随访频率至每月一次",
            "IF 患者连续2次未按时随访 THEN 护士电话提醒，仍未就诊则通知主治医生",
        ]),
    ])

    drugs = [("阿托伐他汀","调脂","20-40mg","口服","肝功能异常"),
             ("二甲双胍","降糖","500-2000mg","口服","eGFR<30"),
             ("赖诺普利","降压","5-40mg","口服","双侧肾动脉狭窄"),
             ("美托洛尔","降压/心率","25-200mg","口服","哮喘/COPD"),
             ("华法林","抗凝","INR指导","口服","活动性出血"),
             ("阿司匹林","抗血小板","100mg","口服","活动性溃疡"),
             ("胰岛素","降糖","个体化","皮下注射","低血糖史"),
             ("氨氯地平","降压","5-10mg","口服","严重肝病"),
             ("呋塞米","利尿","20-80mg","口服/静注","低钾血症"),
             ("地高辛","强心","0.0625-0.25mg","口服","室颤/WPW")]
    drug_rows = [[d[0],d[1],d[2],d[3],d[4],
                  f"{'肾功能' if '肾' in d[4] else '肝功能' if '肝' in d[4] else '血常规'}监测"] for d in drugs]
    for i in range(len(drugs)+1, 41):
        drug_rows.append([f"药物{i:02d}",
                          random.choice(["抗生素","镇痛","抗过敏","止吐","激素"]),
                          f"{random.randint(1,500)}mg",
                          random.choice(["口服","静注","肌注","吸入"]),
                          random.choice(["肾功能不全","肝功能不全","妊娠","过敏史","儿童"]),
                          "用药前监测"])

    patient_rows = []
    diagnoses = ["高血压2级","2型糖尿病","冠心病","COPD","慢性肾病3期","心力衰竭"]
    for i in range(1,81):
        age = random.randint(35,85)
        dx = random.choice(diagnoses)
        hba1c = round(random.uniform(5.5,11.0),1) if "糖尿病" in dx else None
        sbp = random.randint(120,190)
        patient_rows.append([f"P{i:05d}", age, random.choice(["男","女"]),
                             dx, sbp, hba1c, random.choice(["正常","预警","高风险"]),
                             f"2026-{random.randint(1,5):02d}-{random.randint(1,28):02d}",
                             random.choice(["定期随访","失访","已出院","急诊"])])

    write_xlsx(D / "clinical_data.xlsx", [
        ("药品目录", ["药品名","类别","剂量范围","用药途径","主要禁忌","监测项目"], drug_rows),
        ("患者档案", ["患者ID","年龄","性别","主诊断","收缩压","HbA1c%","风险等级","末次就诊","状态"], patient_rows),
        ("不良事件",
         ["事件ID","患者ID","事件类型","严重程度","涉及药物","处理措施","上报状态"],
         [
            ["AE001","P00023","药物过敏","中度","青霉素","停药+抗过敏处理","已上报"],
            ["AE002","P00045","低血糖","重度","胰岛素","葡萄糖静推，ICU观察","已上报"],
            ["AE003","P00067","消化道出血","重度","阿司匹林+华法林","停药+内镜止血","已上报"],
            ["AE004","P00012","肝功能损害","轻度","阿托伐他汀","停药观察","待上报"],
            ["AE005","P00089","肾功能恶化","中度","NSAID","停药+补液","已上报"],
         ]),
    ])

    followup_rows = []
    for i in range(1,101):
        on_time = random.random() > 0.25
        followup_rows.append([f"P{random.randint(1,80):05d}",
                               random.choice(["高血压","糖尿病","冠心病"]),
                               f"2026-{random.randint(1,5):02d}-{random.randint(1,28):02d}",
                               "按时" if on_time else "延误" if random.random()>0.4 else "失访",
                               random.randint(100,190),
                               round(random.uniform(5.5,11.0),1),
                               random.choice(["继续原方案","调整用药","转诊","住院"])])
    write_csv(D / "followup_records.csv",
              ["患者ID","疾病","随访日期","随访状态","血压(SBP)","HbA1c","处理建议"], followup_rows)

    ae_rows = []
    for i in range(1,61):
        ae_rows.append([f"AE{i:04d}", f"P{random.randint(1,80):05d}",
                        random.choice(["药物不良反应","操作并发症","院内感染","跌倒","压疮"]),
                        random.choice(["轻度","中度","重度","死亡"]),
                        random.choice(["是","否"]),
                        f"2026-{random.randint(1,5):02d}-{random.randint(1,28):02d}",
                        random.choice(["已处理","处理中","待评估"])])
    write_csv(D / "adverse_events.csv",
              ["事件ID","患者ID","事件类型","严重程度","是否上报","发生日期","处理状态"], ae_rows)

    write_pptx(D / "clinical_review.pptx", "临床质量月度评审", [
        ("一、患者安全指标", ["院内感染率：1.2‰（目标<2‰）✓","药物不良事件发生率：0.8%","高警示药品错误：0例","30天再入院率：8.3%（目标<10%）✓"]),
        ("二、慢病管理", ["高血压患者血压达标率：72%（目标>70%）✓","糖尿病HbA1c<7%达标率：58%（目标>60%）—未达标","糖尿病患者随访按时率：75%"]),
        ("三、重点预警", ["失访患者：23例（其中高风险12例，已启动外呼）","HbA1c>9%患者：18例，均已启动胰岛素治疗评估","5例多重耐药菌患者已隔离处置"]),
        ("四、质量改进", ["启动糖尿病达标率提升专项：增加营养师联合门诊","优化不良事件上报流程：引入电子化实时上报系统","开展手卫生专项检查：合规率从85%提升至93%"]),
    ])

    write_pdf(D / "drug_safety_report.pdf", "Drug Safety & Pharmacovigilance Report", [
        ("1. Adverse Event Summary Q1 2026", [
            "Total adverse events reported: 47 cases.",
            "Serious adverse events: 8 cases (all reported to regulatory authority within 15 days).",
            "Most common: GI reactions (12 cases), hypoglycemia (8 cases), allergic reactions (7 cases).",
        ]),
        ("2. High-Risk Drug Monitoring", [
            "Warfarin: 45 patients on therapy, 38 (84%) within therapeutic INR range 2.0-3.0.",
            "Insulin: 78 patients, 3 severe hypoglycemia events recorded.",
            "Digoxin: 12 patients, all with serum levels within range 0.5-2.0 ng/mL.",
            "IF INR > 4.0 THEN hold warfarin, notify physician, consider Vitamin K.",
        ]),
        ("3. Drug-Drug Interaction Alerts", [
            "System generated 234 DDI alerts in Q1; 89 overridden by physicians.",
            "Override rate >40% triggers pharmacy review.",
            "Critical DDI (never override): aminoglycosides + vancomycin, warfarin + NSAIDs.",
        ]),
    ])
    print("  医疗: 7个文件生成完成")


# ══════════════════════════════════════════════════════════════════════════════
# 6. 法律
# ══════════════════════════════════════════════════════════════════════════════
def gen_legal():
    D = BASE / "法律"
    print("\n[法律]")
    import random; random.seed(50)

    write_md(D / "legal_framework.md", """# 法务合规管理体系

## 1. 合同管理规则

### 1.1 合同审查标准
| 合同类型 | 金额阈值 | 审查人 | 审查时限 |
|--------|--------|-------|--------|
| 采购合同 | 任意金额 | 法务专员 | 3个工作日 |
| 销售合同 | > 50万 | 法务总监 | 5个工作日 |
| 战略合作协议 | 任意 | 法务总监 + 外部律所 | 10个工作日 |
| 劳动合同 | 无 | HR法务 | 2个工作日 |

### 1.2 合同生效规则
- IF 合同未经法务审核 THEN 视为无效，不得对外签署
- IF 合同金额 > 100万 THEN 须经CEO签署（VP无代签权）
- IF 合同含有排他性条款 THEN 须法务总监和CEO联合审批
- IF 合同内容变更 THEN 须重新法务审核，原版本签署的补充协议须注明"以本协议为准"

### 1.3 合同存档规则
- 所有合同须在签署后3个工作日内上传至合同管理系统（CMS）
- IF 合同金额 > 50万 THEN 纸质原件存档于法务部保险柜，电子版加密存储
- 合同档案保存期限：合同有效期届满后10年（法律规定年限更长者从其规定）

## 2. 知识产权保护规则

### 2.1 商业秘密保护
- 所有员工入职须签署《保密协议》，期限为在职期间+离职后2年
- IF 员工接触核心技术信息 THEN 须额外签署《竞业限制协议》（补偿月薪30%以上）
- IF 竞业补偿 < 月薪30% THEN 竞业限制条款无效
- IF 离职员工加入竞争对手（竞业限制期内）THEN 立即停止支付竞业补偿，评估诉讼可行性

### 2.2 专利管理规则
- IF 员工研发成果满足专利条件 THEN 公司须在发现后6个月内决定是否申请专利
- IF 公司决定不申请专利 THEN 员工可自行申请，但归公司所有
- 专利申请须由研发部门 + 法务部联合评估，外部专利代理机构协助撰写

## 3. 争议解决规则

### 3.1 争议预防规则
- 合同须明确：适用法律（中国大陆法律）、争议解决方式（优先协商→调解→仲裁）
- IF 争议金额 < 50万 THEN 内部法务处理；IF 争议金额 ≥ 50万 THEN 委托外部律所
- 诉讼时效管理：普通合同纠纷3年，劳动争议1年（需要注意中断/中止情形）

### 3.2 仲裁/诉讼启动规则
- IF 协商超过30天未达成协议 THEN 启动仲裁申请
- IF 对方当事人资产有转移风险 THEN 同步申请财产保全
- IF 涉及刑事犯罪 THEN 立即向公安机关报案，同步启动民事追偿

## 4. 合规管理规则

### 4.1 反腐败合规
- 禁止向政府官员、国有企业人员提供超过200元/次的礼品或娱乐
- 所有招待费须事前申请，事后凭发票报销，并记录受邀方信息
- IF 举报合规违规 THEN 举报人受到匿名保护，不得打击报复

### 4.2 数据合规（PIPL/GDPR）
- 个人信息收集须遵循"最小必要"原则，并取得用户授权
- IF 发生个人信息泄露 THEN 须在72小时内向监管机构报告（符合PIPL/GDPR）
- IF 用户行使删除权 THEN 须在30天内完成数据删除，并书面确认
""")

    write_docx(D / "contract_management.docx", "合同管理操作规程", [
        ("1. 合同起草规范", [
            "使用公司标准合同模板；若对方提供格式合同，须逐条与标准条款比对，偏差须法务审核。",
            "必备条款清单：主体信息（全称+统一社会信用代码）、标的、金额、支付方式、交付条件、违约责任、争议解决、保密条款。",
            "禁止条款：单方解除权（需法务特批）、无限连带担保、排除公司基本法律权利的条款。",
        ]),
        ("2. 审核流程", [
            "起草人提交至合同管理系统（CMS）→ 法务专员初审（1个工作日）→ 法务总监终审（2个工作日）→ 业务VP确认 → 签署",
            "重大合同（>100万/含排他/含知识产权转让）：增加外部律所审核环节（5个工作日）",
            "IF 法务提出实质性修改意见 THEN 业务部门须与法务共同讨论后决定是否接受，不得单方面覆盖法务意见",
        ]),
        ("3. 合同履行监控", [
            "法务系统自动提醒：合同到期前90/60/30天分别发送提醒至合同负责人和法务",
            "付款里程碑到期3天内法务发送催款提醒；超期7天启动正式催收",
            "IF 合同标的物未按时交付 THEN 通知对方书面确认，超期15天触发违约责任条款",
        ]),
        ("4. 争议处理", [
            "发现合同纠纷第一时间通知法务，禁止业务自行承诺或妥协",
            "保留所有往来邮件、微信记录、会议纪要等证据",
            "IF 对方提出诉讼/仲裁 THEN 法务须在收到文书后48小时内评估并制定应对策略",
        ]),
    ])

    contract_rows = []
    types = ["销售合同","采购合同","劳动合同","服务协议","战略合作","保密协议","租赁合同"]
    statuses = ["有效","到期预警","已到期","谈判中","争议中"]
    for i in range(1,61):
        amt = random.choice([random.randint(1,50), random.randint(50,500), random.randint(500,5000)])*10000
        contract_rows.append([f"CT-2026-{i:04d}", random.choice(types),
                               f"甲方{i:03d}有限公司", f"乙方{i:03d}有限公司",
                               amt, f"2025-{random.randint(1,12):02d}-01",
                               f"2027-{random.randint(1,12):02d}-01",
                               random.choice(statuses),
                               random.choice(["北京仲裁委","上海仲裁委","人民法院"])])

    dispute_rows = []
    for i in range(1,31):
        amt = random.randint(10,500)*10000
        dispute_rows.append([f"DIS-2026-{i:04d}",
                              random.choice(["合同违约","货款纠纷","知识产权","劳动争议","侵权"]),
                              amt,
                              random.choice(["协商中","仲裁","诉讼","已结案"]),
                              f"2025-{random.randint(1,12):02d}-{random.randint(1,28):02d}",
                              random.choice(["内部法务","外部律所"]),
                              random.choice(["胜诉","败诉","和解","进行中"])])

    write_xlsx(D / "legal_database.xlsx", [
        ("合同台账", ["合同编号","合同类型","甲方","乙方","金额(元)","签署日期","到期日期","状态","争议解决"], contract_rows),
        ("争议案件", ["案件编号","争议类型","涉案金额","处理阶段","立案日期","承办方","结果"], dispute_rows),
        ("合规检查",
         ["检查日期","检查项目","部门","问题描述","风险等级","整改期限","整改状态"],
         [
            ["2026-03-15","反腐败培训覆盖率","全公司","覆盖率92%（目标95%）","低","2026-04-30","已整改"],
            ["2026-03-15","合同审核合规","销售部","2份合同未经法务审核","高","2026-03-30","已整改"],
            ["2026-04-10","数据合规PIPL","产品研发","隐私政策未更新","中","2026-05-10","整改中"],
            ["2026-04-10","劳动合同","HR","3名实习生未签劳动协议","高","2026-04-20","已整改"],
         ]),
    ])

    ip_rows = []
    for i in range(1,61):
        ip_rows.append([f"IP-{i:04d}",
                        random.choice(["发明专利","实用新型","商标","著作权","商业秘密"]),
                        f"知识产权名称{i:03d}",
                        random.choice(["已申请","已授权","申请中","已放弃"]),
                        f"2020-{random.randint(1,12):02d}-01",
                        f"2030-{random.randint(1,12):02d}-01",
                        random.choice(["核心技术","品牌保护","商业方法","软件著作权"])])
    write_csv(D / "ip_portfolio.csv",
              ["IP编号","类型","名称","状态","申请日期","有效期至","保护领域"], ip_rows)

    regulation_rows = []
    for i in range(1,41):
        regulation_rows.append([f"REG-{i:04d}",
                                 random.choice(["PIPL","GDPR","劳动法","公司法","税法","反垄断法","网络安全法"]),
                                 f"2026-{random.randint(1,6):02d}-{random.randint(1,28):02d}",
                                 random.choice(["新法规","修订","司法解释","合规指引"]),
                                 random.choice(["高","中","低"]),
                                 random.choice(["已评估","评估中","待评估"])])
    write_csv(D / "regulatory_changes.csv",
              ["记录ID","法规名称","生效日期","变更类型","影响等级","应对状态"], regulation_rows)

    write_pptx(D / "legal_review.pptx", "法务季度工作汇报", [
        ("一、合同管理概况", ["本季度新签合同：186份（总金额2.3亿）","合同纠纷率：1.6%（行业均值3%）✓","平均审核周期：2.8个工作日","重大合同（>100万）：15份，全部合规"]),
        ("二、争议处理进展", ["在审案件：12件（合同违约8件，劳动争议4件）","胜/和解率：83%","本季度新增2件，结案3件","最大标的：货款纠纷350万（仲裁中）"]),
        ("三、合规风险", ["高风险：PIPL数据合规（产品团队隐私政策未更新）—已催促","中风险：2名离职员工可能违反竞业限制—已发律师函","低风险：日常合规培训覆盖率达96%"]),
        ("四、知识产权", ["本季度申请专利：8件","商标注册：完成3个新品牌注册","竞争对手侵权：发现1起，律师函已发送"]),
        ("五、下季度重点", ["完成PIPL合规整改","启动全员反腐培训（线上+测试）","推进核心技术专利申请（5件）"]),
    ])

    write_pdf(D / "compliance_manual.pdf", "Legal Compliance Manual", [
        ("1. Anti-Corruption Policy", [
            "Prohibited: Gifts >200 RMB per occasion to government officials or SOE employees.",
            "All entertainment expenses require pre-approval and post-receipt documentation.",
            "IF gift/entertainment >200 RMB THEN prior written approval from Legal required.",
            "Whistleblower hotline available 24/7; anonymous reports protected.",
        ]),
        ("2. Data Privacy (PIPL/GDPR)", [
            "Minimum necessary principle: collect only data required for stated purpose.",
            "User consent must be explicit, specific, and documented.",
            "IF data breach detected THEN notify regulator within 72 hours.",
            "IF user requests data deletion THEN fulfill within 30 days and confirm in writing.",
        ]),
        ("3. Contract Compliance", [
            "No contract may be signed without Legal review.",
            "Contracts >1M RMB require CEO signature; no delegation allowed.",
            "All contracts uploaded to CMS within 3 business days of signing.",
            "IF contract terms deviate from standard THEN Legal must document rationale.",
        ]),
    ])
    print("  法律: 7个文件生成完成")


# ══════════════════════════════════════════════════════════════════════════════
# 7. 教育
# ══════════════════════════════════════════════════════════════════════════════
def gen_education():
    D = BASE / "教育"
    print("\n[教育]")
    import random; random.seed(60)

    write_md(D / "academic_policy.md", """# 学术管理制度

## 1. 学分与毕业规则

### 1.1 学分体系
| 课程类型 | 最低学分 | 说明 |
|--------|--------|-----|
| 必修课 | 64学分 | 含核心基础课32学分+专业必修32学分 |
| 选修课 | 32学分 | 含专业选修24学分+公选课8学分 |
| 实践学分 | 16学分 | 实习、毕业论文、社会实践 |
| 总计 | 128学分 | 须全部完成方可申请毕业 |

### 1.2 毕业资格规则
- IF 修满128学分 AND GPA ≥ 2.0 AND 毕业论文通过 THEN 授予学士学位
- IF GPA < 2.0 AND 学分达标 THEN 仅颁发毕业证书，不授予学士学位
- IF 必修课有挂科未补考 THEN 不允许申请毕业
- IF 学习年限 > 6年（弹性学制） THEN 自动退学，不得延期

### 1.3 课程解锁规则
- IF 前置课程成绩 < 60分 THEN 不允许注册对应高阶课程
- IF 前置课程未修完 THEN 系统自动拒绝注册
- IF 重修次数 > 2次 THEN 须向教务处申请特殊许可

## 2. 学业预警规则

### 2.1 预警触发条件
- IF 单科成绩 < 60分 THEN 一级预警：辅导员发送提醒，建议学生到教授答疑
- IF GPA < 2.0 连续2个学期 THEN 二级预警：辅导员 + 家长通知，制定学业改进计划
- IF GPA < 1.5 THEN 三级预警：启动留级评估程序，由教务委员会审查
- IF 旷课超过课程总时数1/3 THEN 取消该课程考试资格，成绩记为0分

### 2.2 考试资格规则
- IF 出勤率 < 70% THEN 不允许参加期末考试（特殊情况须教务处审批）
- IF 作业提交率 < 60% THEN 课程成绩上限为70分
- IF 作弊行为查实 THEN 该课程成绩记为0分，且计入学籍档案

## 3. 奖学金评定规则

### 3.1 评定标准
| 奖学金类型 | GPA要求 | 其他条件 |
|---------|--------|--------|
| 一等奖学金 | ≥3.8 | 无挂科，综合素质A等 |
| 二等奖学金 | ≥3.5 | 无挂科 |
| 三等奖学金 | ≥3.2 | 当学期无不及格 |
| 国家奖学金 | ≥3.9 | 排名前2%，品德优秀 |

### 3.2 奖学金撤销规则
- IF 获奖后发现学术不端（抄袭/作弊）THEN 撤销当年奖学金并通报批评
- IF 获奖后当学期GPA < 2.0 THEN 追回当学期奖学金
- IF 违纪记录（警告及以上）THEN 当学年不得申请奖学金

## 4. 课程管理规则
- 退课截止日：每学期第4周末（超过不得退课，成绩记为F）
- 缓考申请：须在考前3天提交，仅接受病假（须附医院证明）
- IF 缓考申请未获批准 THEN 缺考记为0分
- IF 开课人数 < 15人 THEN 该课程自动取消，学生转入备选课程
""")

    write_docx(D / "curriculum_design.docx", "课程设计与教学规范", [
        ("1. 课程体系框架", [
            "第一学年：通识基础课（高等数学、线性代数、大学英语、程序设计基础）",
            "第二学年：专业基础课（数据结构、算法分析、数据库原理、操作系统）",
            "第三学年：专业核心课（机器学习、深度学习、计算机视觉、NLP）+ 方向选修",
            "第四学年：毕业设计（30学分）+ 企业实习（6学分）",
        ]),
        ("2. 课程前置要求", [
            "机器学习 → 前置：高等数学（微积分）≥75分 + 线性代数≥70分 + Python基础≥70分",
            "深度学习 → 前置：机器学习≥75分",
            "计算机视觉 → 前置：深度学习≥70分",
            "自然语言处理 → 前置：深度学习≥70分 + 概率论≥65分",
        ]),
        ("3. 教学质量规则", [
            "课程满意度调查每学期进行，低于3.5/5.0触发教学督导介入",
            "IF 期末成绩分布：优秀率（≥90分）> 30% THEN 审查评分标准是否过于宽松",
            "IF 不及格率 > 30% THEN 审查教学难度和题目合理性",
            "助教授课不得超过总课时的20%，超过须经教学委员会审批",
        ]),
        ("4. 毕业论文管理", [
            "选题确认：大四第一学期第8周前完成选题审批",
            "开题报告：第10周前提交，导师+教研室主任联合审核",
            "中期检查：第二学期第6周，导师检查研究进度，进度<30%发出预警",
            "答辩资格：IF 查重率 > 15% THEN 不允许参加答辩，须修改后重新提交",
            "答辩结果：三票以上委员同意方可通过（答辩委员会5人）",
        ]),
    ])

    courses = [("CS101","高等数学I","必修",4,None,85),
               ("CS102","线性代数","必修",3,None,78),
               ("CS201","数据结构","必修",4,"CS101",82),
               ("CS202","算法分析","必修",4,"CS201",75),
               ("CS301","机器学习","专业核心",4,"CS201,CS102",68),
               ("CS302","深度学习","专业核心",4,"CS301",72),
               ("CS303","计算机视觉","专业选修",3,"CS302",65),
               ("CS304","自然语言处理","专业选修",3,"CS302",70),
               ("CS401","毕业设计","实践",10,None,None),
               ("GE101","大学英语","必修",4,None,None)]
    course_rows = [[c[0],c[1],c[2],c[3],c[4] or "无",f"{c[5]}分" if c[5] else "面评",
                    random.randint(30,120), random.choice(["开课","暂停","已结课"])] for c in courses]
    for i in range(11,31):
        course_rows.append([f"CS{i:03d}", f"课程{i:03d}", random.choice(["必修","选修","实践"]),
                            random.randint(2,4), random.choice(["CS101","CS201","无"]),
                            f"{random.randint(60,90)}分", random.randint(20,80),"开课"])

    student_rows = []
    for i in range(1,101):
        gpa = round(random.uniform(1.2,4.0),2)
        credit = random.randint(60,128)
        student_rows.append([f"STU{i:05d}", f"学生{i:03d}", random.randint(2022,2025),
                             random.choice(["计算机科学","软件工程","人工智能","数据科学"]),
                             gpa, credit,
                             "正常" if gpa>=2.0 else "二级预警" if gpa>=1.5 else "三级预警",
                             random.choice(["在读","休学","毕业","退学"])])

    write_xlsx(D / "academic_data.xlsx", [
        ("课程目录", ["课程编号","课程名称","类型","学分","前置课程","通过标准","当前选课人数","状态"], course_rows),
        ("学生档案", ["学号","姓名","入学年份","专业","GPA","已修学分","预警状态","学籍状态"], student_rows),
        ("奖学金记录",
         ["学号","奖学金类型","金额(元)","评定学期","GPA","综合排名","是否通过"],
         [
            ["STU00001","国家奖学金",8000,"2025-秋",3.92,"前2%","是"],
            ["STU00005","一等奖学金",3000,"2025-秋",3.85,"前5%","是"],
            ["STU00012","二等奖学金",2000,"2025-秋",3.55,"前15%","是"],
            ["STU00023","三等奖学金",1000,"2025-秋",3.25,"前30%","是"],
            ["STU00034","申请一等","—","2025-秋",3.6,"—","否（有警告记录）"],
         ]),
    ])

    grade_rows = []
    for i in range(1,121):
        score = random.randint(35,100)
        grade_rows.append([f"STU{random.randint(1,100):05d}",
                           random.choice(["CS101","CS201","CS301","CS302","CS303"]),
                           f"2025-秋",
                           score,
                           "优秀" if score>=90 else "良好" if score>=80 else "中等" if score>=70 else "及格" if score>=60 else "不及格",
                           random.choice(["正常","缺考","作弊","缓考"])])
    write_csv(D / "grade_records.csv",
              ["学号","课程编号","学期","成绩","等级","备注"], grade_rows)

    attendance_rows = []
    for i in range(1,101):
        rate = round(random.uniform(0.5,1.0),2)
        attendance_rows.append([f"STU{i:05d}", random.choice(["CS301","CS302","CS303"]),
                                 f"2025-秋", f"{rate:.0%}",
                                 "正常" if rate>=0.7 else "预警",
                                 random.randint(0,8)])
    write_csv(D / "attendance.csv",
              ["学号","课程","学期","出勤率","状态","缺课次数"], attendance_rows)

    write_pptx(D / "academic_review.pptx", "学院学术质量年度报告", [
        ("一、学生规模", ["在读学生：2,380人","本年度新生：620人","毕业生：580人","留级率：3.2%（行业均值4%）✓"]),
        ("二、学业成绩", ["平均GPA：3.12","GPA≥3.5优秀率：22%","不及格率：8.5%（目标<10%）✓","挂科重修率：15%"]),
        ("三、预警情况", ["一级预警：156人（已全部约谈）","二级预警：34人（含PIP计划）","三级预警：8人（留级评估中）"]),
        ("四、毕业论文", ["答辩通过率：94%","优秀论文：45篇（占比7.7%）","查重超标返修：23篇","平均查重率：8.3%（目标<15%）✓"]),
        ("五、质量改进", ["引入AI辅助学业预警系统","推进双导师制（校内+企业）","扩大国际交流项目名额至80人"]),
    ])

    write_pdf(D / "accreditation_report.pdf", "Program Accreditation Self-Assessment", [
        ("1. Program Overview", [
            "Program: B.Sc. in Artificial Intelligence, established 2020.",
            "Current enrollment: 2,380 students; Faculty: 85 full-time, 20 adjunct.",
            "Graduation rate (4-year): 88%; Employment rate: 96%.",
        ]),
        ("2. Learning Outcomes Assessment", [
            "LO1 (Algorithmic Thinking): 85% of graduates meet target competency.",
            "LO2 (Programming Skills): 91% proficiency in Python/C++.",
            "LO3 (AI/ML Foundations): 78% meet advanced competency (target: 80%) — improvement needed.",
            "LO4 (Ethical AI): 72% pass ethics assessment — curriculum enhancement planned.",
        ]),
        ("3. Continuous Improvement", [
            "Student feedback: 3.8/5.0 average; lowest scores for course load management.",
            "Faculty development: 100% completed pedagogy training.",
            "IF course satisfaction < 3.5 THEN curriculum revision required within 1 semester.",
            "Industry advisory board meets quarterly; curriculum updated annually.",
        ]),
    ])
    print("  教育: 7个文件生成完成")


# ══════════════════════════════════════════════════════════════════════════════
# 主程序
# ══════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("=" * 60)
    print("多领域测试数据生成器")
    print("=" * 60)
    generators = [gen_supply_chain, gen_hr, gen_finance, gen_marketing, gen_medical, gen_legal, gen_education]
    ok, fail = 0, []
    for gen in generators:
        try:
            gen(); ok += 1
        except Exception as e:
            import traceback
            print(f"  ❌ 失败: {e}")
            traceback.print_exc()
            fail.append(gen.__name__)
    print(f"\n{'='*60}")
    print(f"完成: {ok}/{len(generators)} 个领域")
    if fail: print(f"失败: {fail}")
    print(f"文件保存至: {BASE}")
